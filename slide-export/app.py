"""slide-export — Slide Studio HTML 덱을 16:9 PDF / PPTX 로 렌더.

POST /pdf   {"html": "…<section class='slide'>…"}  → application/pdf
POST /pptx  {"html": …}                            → .pptx
GET  /health

두 형식 모두 **슬라이드 단위 캡처** 방식: 덱의 각 `.slide` 만 차례로 active
토글 → headless chromium 으로 1280×720 PNG 캡처 → PDF(Pillow)·PPTX(python-pptx)로
슬라이드당 이미지 1장 조립. N 슬라이드 = N 페이지 보장(CSS page-break
누락 이슈 없음). 텍스트 편집성 상실, 단 그라데이션/SVG차트/한글 디자인은 화면 그대로.
"""
from __future__ import annotations

import contextlib
from io import BytesIO

from fastapi import FastAPI, HTTPException, Request, Response
from PIL import Image
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

# 슬라이드 i 만 active 토글 (덱의 go() 의존 없이 직접 — 견고).
_ACTIVATE = (
    "(k)=>{document.querySelectorAll('.slide')"
    ".forEach((s,i)=>s.classList.toggle('active', i===k))}"
)
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
# 1280×720 px @ 96dpi = 13.333×7.5 in = 16:9 (PDF/PPTX 동일 페이지 크기).
_DPI = 96.0

_state: dict = {}


@contextlib.asynccontextmanager
async def lifespan(_app: FastAPI):
    pw = await async_playwright().start()
    _state["pw"] = pw
    _state["browser"] = await pw.chromium.launch(args=["--no-sandbox"])
    try:
        yield
    finally:
        await _state["browser"].close()
        await pw.stop()


app = FastAPI(lifespan=lifespan, title="slide-export")


@app.get("/health")
async def health() -> dict:
    return {"status": "ok"}


async def _slide_shots(html: str) -> list[bytes]:
    """덱의 각 .slide 차례로 active 토글 → 1280×720 PNG 캡처."""
    if not isinstance(html, str) or not html.strip():
        raise HTTPException(400, "missing 'html'")
    page = await _state["browser"].new_page(viewport={"width": 1280, "height": 720})
    try:
        await page.set_content(html, wait_until="load")
        deck = page.locator(".deck")
        n = await page.locator(".slide").count()
        if n == 0:
            raise HTTPException(422, "no .slide elements")
        shots: list[bytes] = []
        for k in range(n):
            await page.evaluate(_ACTIVATE, k)
            await page.wait_for_timeout(120)
            shots.append(await deck.screenshot(type="png"))
        return shots
    finally:
        await page.close()


async def _html_from(request: Request) -> str:
    return (await request.json()).get("html")


@app.post("/pdf")
async def to_pdf(request: Request) -> Response:
    """슬라이드당 이미지 1페이지의 16:9 PDF."""
    shots = await _slide_shots(await _html_from(request))
    imgs = [Image.open(BytesIO(p)).convert("RGB") for p in shots]
    buf = BytesIO()
    imgs[0].save(buf, format="PDF", save_all=True,
                 append_images=imgs[1:], resolution=_DPI)
    return Response(
        content=buf.getvalue(),
        media_type="application/pdf",
        headers={"Content-Disposition": 'attachment; filename="presentation.pdf"'},
    )


@app.post("/pptx")
async def to_pptx(request: Request) -> Response:
    """슬라이드당 이미지 1장의 16:9 PPTX."""
    shots = await _slide_shots(await _html_from(request))
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in shots:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(BytesIO(png), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    buf = BytesIO()
    prs.save(buf)
    return Response(
        content=buf.getvalue(),
        media_type=_PPTX_MIME,
        headers={"Content-Disposition": 'attachment; filename="presentation.pptx"'},
    )
