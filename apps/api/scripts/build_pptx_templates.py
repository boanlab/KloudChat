"""Generates `template.pptx` (master and layouts, no slides) and `form.pptx` per 덱 서식.

Run from the API image:

    docker compose run --rm --no-deps -v "$PWD/apps/api:/repo" -w /repo api \
        python scripts/build_pptx_templates.py

Design lives on the master; every line of text goes into a layout placeholder.
Word counterpart: `build_docx_templates.py`. Fonts must be faces every Hangul
Windows and Hancom install carries, or PowerPoint substitutes silently.
"""

from __future__ import annotations

import pathlib
import sys

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

ROOT = pathlib.Path(__file__).resolve().parent.parent / "app" / "design_templates"

#: 16:9 slide size in EMU.
WIDE = (Emu(12192000), Emu(6858000))


class Slide:
    """One slide of the blank form, made from the named layout."""

    def __init__(self, layout: str, title: str, *body: str, second: tuple[str, ...] = ()) -> None:
        self.layout = layout
        self.title = title
        self.body = body
        self.second = second


class Spec:
    """One 덱 서식's PowerPoint half."""

    def __init__(
        self,
        folder: str,
        *,
        body_font: str,
        heading_font: str,
        accent: str,
        ink: str,
        paper: str,
        title_pt: float,
        heading_pt: float,
        body_pt: float,
        note: str,
        form: tuple[Slide, ...] = (),
    ) -> None:
        self.folder = folder
        self.body_font = body_font
        self.heading_font = heading_font
        self.accent = accent
        self.ink = ink
        self.paper = paper
        self.title_pt = title_pt
        self.heading_pt = heading_pt
        self.body_pt = body_pt
        self.note = note
        self.form = form


#: Mirrors each 서식's `seed.html` type scale; keep the two in step.
_SPECS = (
    Spec(
        "deck-lecture",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="1F6FEB",
        ink="1A1A1A",
        paper="FFFFFF",
        title_pt=40,
        heading_pt=30,
        body_pt=20,
        note="강의형. 뒤에서도 읽히게 크게.",
        form=(
            Slide("Title Slide", "강의 제목", "부제 · 날짜 · 강사"),
            Slide(
                "Title and Content",
                "오늘 다루는 것",
                "한 장에 개념 하나.",
                "글머리는 다섯 줄을 넘기지 않습니다.",
                "말로 할 것은 슬라이드에 적지 않습니다.",
            ),
            Slide("Section Header", "1. 첫 번째 주제", "이 절에서 답할 질문 한 줄"),
            Slide(
                "Two Content",
                "개념과 예시",
                "개념: 한 문장으로 정의합니다.",
                second=("예시: 그 정의가 들어맞는 구체적인 경우 하나.",),
            ),
            Slide(
                "Title and Content",
                "정리",
                "오늘 남길 문장 하나.",
                "다음 시간에 이어지는 것.",
            ),
        ),
    ),
    Spec(
        "deck-proposal",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="0F766E",
        ink="1A1A1A",
        paper="FFFFFF",
        title_pt=36,
        heading_pt=28,
        body_pt=18,
        note="제안. 문제·접근·근거를 차례로 놓는다.",
        form=(
            Slide("Title Slide", "제안 제목", "고객사 · 제안일 · 담당"),
            Slide(
                "Title and Content",
                "고객의 과제",
                "지금 무엇이 문제인지 고객의 말로 적습니다.",
                "그대로 두면 무엇이 되는지 숫자로.",
            ),
            Slide(
                "Title and Content",
                "제안하는 것",
                "무엇을 하겠다는 것인지 한 문장.",
                "그것이 위 과제의 어느 부분을 푸는지.",
            ),
            Slide(
                "Comparison",
                "지금과 이후",
                "지금",
                "현재 방식으로 드는 시간·비용",
                second=("도입 후", "달라지는 수치와 그 근거"),
            ),
            Slide(
                "Title and Content",
                "도입 일정",
                "단계와 기간을 적습니다.",
                "고객이 준비해야 할 것을 함께 적습니다.",
            ),
            Slide("Title and Content", "요청", "이 자리에서 무엇을 결정해 달라는 것인지."),
        ),
    ),
    Spec(
        "deck-editorial",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="B45309",
        ink="1A1A1A",
        paper="FFFFFF",
        title_pt=36,
        heading_pt=26,
        body_pt=17,
        note="편집형. 제목 아래 룰 하나, 본문은 왼쪽.",
        form=(
            Slide("Title Slide", "제목", "부제 한 줄"),
            Slide("Section Header", "장 제목", "이 장이 답하는 질문"),
            Slide(
                "Content with Caption",
                "본문 제목",
                "본문. 한 화면에 한 가지만 말합니다.",
                second=("옆에 붙는 짧은 설명이나 출처.",),
            ),
            Slide(
                "Picture with Caption",
                "그림이 들어가는 장",
                "그림을 여기에 넣습니다.",
                second=("그림이 무엇을 보여 주는지 한 줄.",),
            ),
            Slide("Title and Content", "맺음", "남길 문장 하나."),
        ),
    ),
    Spec(
        "deck-signal",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="F59E0B",
        ink="F5F5F5",
        paper="111111",
        title_pt=44,
        heading_pt=32,
        body_pt=22,
        note="대비형. 어두운 화면에 큰 글씨.",
        form=(
            Slide("Title Slide", "제목", "부제 한 줄"),
            Slide("Title Only", "한 화면에 한 문장"),
            Slide(
                "Title and Content",
                "숫자로 말하는 장",
                "수치 하나와 그 뜻 한 줄.",
                "출처는 작게, 아래에.",
            ),
            Slide("Section Header", "전환", "여기서부터 무엇이 달라지는지"),
            Slide("Title Only", "마지막에 남길 한 문장"),
        ),
    ),
    Spec(
        "deck-case",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="1F6FEB",
        ink="1A1A1A",
        paper="FFFFFF",
        title_pt=36,
        heading_pt=28,
        body_pt=18,
        note="케이스 발표. 대안을 나란히 놓는다.",
        form=(
            Slide("Title Slide", "케이스 제목", "대상 기업 · 발표자 · 날짜"),
            Slide(
                "Title and Content",
                "권고",
                "무엇을 하자는 것인지 한 문장. 뒤의 장이 이것의 근거가 됩니다.",
            ),
            Slide("Title and Content", "현황", "수치 하나와 그 뜻 한 줄.", "출처는 아래에 작게."),
            Slide(
                "Comparison",
                "대안 비교",
                "안 1",
                "같은 기준에서 본 장단점",
                second=("안 2", "같은 기준에서 본 장단점"),
            ),
            Slide(
                "Title and Content", "위험",
                "틀렸을 때 무엇이 일어나는가.", "무엇을 보면 아는가.",
            ),
            Slide("Title and Content", "요청", "무엇을 판단해 달라는 것인지."),
        ),
    ),
    Spec(
        "deck-defense",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="4B3B8F",
        ink="1A1A1A",
        paper="FFFFFF",
        title_pt=34,
        heading_pt=26,
        body_pt=17,
        note="심사 발표. 질문·방법·결과·한계.",
        form=(
            Slide("Title Slide", "연구 제목", "소속 · 이름 · 심사 단계 · 날짜"),
            Slide("Title and Content", "연구 질문", "무엇을 묻는 연구인지 한 문장으로."),
            Slide("Title and Content", "방법", "어떻게 봤는지.", "표본·기간·도구."),
            Slide(
                "Picture with Caption",
                "결과",
                "그림이나 표를 여기에.",
                second=("이 그림이 말하는 것 한 줄.",),
            ),
            Slide("Title and Content", "한계", "무엇을 못 봤는지, 왜 못 봤는지."),
            Slide("Title and Content", "다음 계획", "남은 기간에 무엇을 하는지."),
        ),
    ),
    Spec(
        "deck-briefing",
        body_font="맑은 고딕",
        heading_font="맑은 고딕",
        accent="0F766E",
        ink="1A1A1A",
        paper="FFFFFF",
        title_pt=34,
        heading_pt=26,
        body_pt=18,
        note="사내 브리핑. 정해진 것과 정할 것을 가른다.",
        form=(
            Slide("Title Slide", "브리핑 제목", "일시 · 대상 · 작성"),
            Slide("Title and Content", "오늘 정할 것", "무엇을 결정하러 모였는지 한 줄."),
            Slide(
                "Two Content",
                "정해진 것 / 정할 것",
                "이미 정해진 것.",
                second=("아직 정하지 못한 것.",),
            ),
            Slide(
                "Title and Content",
                "진행 상황",
                "항목마다 무엇까지 됐는지. '진행 중' 은 상태가 아닙니다.",
            ),
            Slide("Title and Content", "막힌 것", "무엇 때문에 막혔고 누가 풀 수 있는지."),
            Slide("Title and Content", "담당과 기한", "누가, 무엇을, 언제까지."),
        ),
    ),
)


_THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


def _theme(presentation):
    """Returns the master's theme part and its parsed XML tree."""
    part = presentation.slide_masters[0].part.part_related_by(_THEME_REL)
    return part, etree.fromstring(part.blob)


def _seal(part, tree) -> None:
    part._blob = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)


def _design(presentation, spec: Spec) -> None:
    """Writes fonts and colours into the theme."""
    part, tree = _theme(presentation)
    elements = tree.find(qn("a:themeElements"))

    # Hangul runs read `a:ea`, not `a:latin`.
    scheme = elements.find(qn("a:fontScheme"))
    for tag, font in (("a:majorFont", spec.heading_font), ("a:minorFont", spec.body_font)):
        block = scheme.find(qn(tag))
        block.find(qn("a:latin")).set("typeface", font)
        block.find(qn("a:ea")).set("typeface", font)

    colours = elements.find(qn("a:clrScheme"))
    for tag, value in (
        ("a:dk1", spec.ink),
        ("a:lt1", spec.paper),
        ("a:accent1", spec.accent),
    ):
        block = colours.find(qn(tag))
        for child in list(block):
            block.remove(child)
        etree.SubElement(block, qn("a:srgbClr")).set("val", value)

    _seal(part, tree)


def _sizes(presentation, spec: Spec) -> None:
    """Sets type sizes on the master's text styles.

    The size must go on `a:defRPr` inside each level; set on `a:lvlNpPr`
    itself it is silently ignored.
    """
    master = presentation.slide_masters[0]
    styles = master.element.find(qn("p:txStyles"))
    if styles is None:
        return
    for tag, size in (
        ("p:titleStyle", spec.title_pt),
        ("p:bodyStyle", spec.body_pt),
        ("p:otherStyle", spec.body_pt),
    ):
        block = styles.find(qn(tag))
        if block is None:
            continue
        for level, properties in enumerate(block):
            # 2pt per outline level, floor 12pt.
            step = max(size - level * 2, 12)
            run = properties.find(qn("a:defRPr"))
            if run is None:
                run = etree.SubElement(properties, qn("a:defRPr"))
            run.set("sz", str(int(step * 100)))


def _widen(presentation) -> None:
    """Scales master and layout placeholders horizontally from python-pptx's 4:3 default to `WIDE`.

    Must run before any slide is added: slides clone their layout's geometry.
    """
    factor = WIDE[0] / Emu(9144000)

    # Read all geometry before writing any: a layout placeholder without its
    # own geometry reports the master's and would otherwise be scaled twice.
    held = [
        (shape, shape.left, shape.top, shape.width, shape.height)
        for holder in [presentation.slide_masters[0], *presentation.slide_layouts]
        for shape in holder.shapes
        if None not in (shape.left, shape.top, shape.width, shape.height)
    ]
    # Write all four: setting one value on a shape without `a:xfrm` creates
    # the element with the other half of the pair at zero.
    for shape, left, top, width, height in held:
        shape.left = Emu(int(left * factor))
        shape.top = Emu(int(top))
        shape.width = Emu(int(width * factor))
        shape.height = Emu(int(height))


def _blank(spec: Spec) -> Presentation:
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = WIDE
    _widen(presentation)
    _design(presentation, spec)
    _sizes(presentation, spec)
    return presentation


def build(spec: Spec) -> pathlib.Path:
    """Writes `template.pptx`: master and layouts, no slides (the writer appends to it)."""
    out = ROOT / spec.folder / "template.pptx"
    _blank(spec).save(out)
    return out


def _layout(presentation, name: str):
    for layout in presentation.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"레이아웃 없음: {name}")


def build_form(spec: Spec) -> pathlib.Path | None:
    """Writes `form.pptx`, the downloadable blank form. Placeholders only, no text boxes."""
    if not spec.form:
        return None
    presentation = _blank(spec)

    for wanted in spec.form:
        layout = _layout(presentation, wanted.layout)
        slide = presentation.slides.add_slide(layout)
        placeholders = {p.placeholder_format.idx: p for p in slide.placeholders}

        if 0 in placeholders:
            placeholders[0].text_frame.text = wanted.title

        body = placeholders.get(1)
        if body is not None and wanted.body:
            frame = body.text_frame
            frame.text = wanted.body[0]
            for line in wanted.body[1:]:
                frame.add_paragraph().text = line
        elif body is not None:
            # PowerPoint drops an empty placeholder on open; a space keeps it.
            body.text_frame.text = " "

        # Second column of `Comparison` / `Two Content` layouts.
        for index, line in zip((2, 3, 4), wanted.second, strict=False):
            if index in placeholders:
                placeholders[index].text_frame.text = line

        # Date, footer and slide-number placeholders are left to the master.

    out = ROOT / spec.folder / "form.pptx"
    presentation.save(out)
    return out


def main() -> int:
    for spec in _SPECS:
        folder = ROOT / spec.folder
        if not folder.is_dir():
            print(f"없는 서식: {spec.folder}", file=sys.stderr)
            return 1
        path = build(spec)
        form = build_form(spec)
        blank = f" + 양식 {form.stat().st_size:,}바이트" if form else ""
        print(f"{spec.folder:15} {path.stat().st_size:>7,}바이트{blank}  {spec.note}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
