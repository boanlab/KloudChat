"""Deck export to `.pptx` and `.pdf`, one 16:9 slide per page from the stored artifact fields.

Both share the 960×540 pt geometry (13.333×7.5 in). Speaker notes go to the
.pptx notes pane only.
"""

from __future__ import annotations

import base64
import io
import logging
import re
from dataclasses import dataclass
from html.parser import HTMLParser

import PIL.Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from reportlab.lib.colors import Color
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen import canvas

from app.services import charts, deck, deck_type, design, fonts, pictures

log = logging.getLogger(__name__)


class _InlineRuns(HTMLParser):
    """Reads the editor's sanitised inline HTML into `(text, style)` runs."""

    def __init__(self) -> None:
        super().__init__()
        self.stack: list[dict] = [{}]
        self.runs: list[tuple[str, dict]] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        style = dict(self.stack[-1])
        values = dict(attrs)
        if tag in ("b", "strong"):
            style["bold"] = True
        if tag in ("i", "em"):
            style["italic"] = True
        if tag == "u":
            style["underline"] = True
        if tag == "font":
            if (size := values.get("size")) and str(size).isdigit():
                style["size"] = int(size)
            if color := values.get("color"):
                style["color"] = color
        if tag == "span" and (css := values.get("style")):
            rules = {
                key.strip().lower(): value.strip()
                for rule in css.split(";")
                if ":" in rule
                for key, value in [rule.split(":", 1)]
            }
            if re.fullmatch(r"[0-9.]+em", rules.get("font-size", "")):
                style["scale"] = float(rules["font-size"][:-2])
            if rules.get("font-weight") in ("bold", "600", "700", "800", "900"):
                style["bold"] = True
            if rules.get("font-style") == "italic":
                style["italic"] = True
            if "underline" in rules.get("text-decoration", ""):
                style["underline"] = True
            if color := rules.get("color"):
                style["color"] = color
        self.stack.append(style)

    def handle_endtag(self, _tag: str) -> None:
        if len(self.stack) > 1:
            self.stack.pop()

    def handle_data(self, data: str) -> None:
        if data:
            self.runs.append((data, dict(self.stack[-1])))


def _inline_runs(html: str | None, fallback: str) -> list[tuple[str, dict]]:
    if not html:
        return [(fallback, {})]
    parser = _InlineRuns()
    parser.feed(html)
    return parser.runs or [(fallback, {})]


#: 16:9 in points (EMU/12700, reportlab's default unit); drives both exporters.
_W, _H = 960.0, 540.0

#: Body-slide geometry from the shared type scale (`deck_type`), in points: the panel's
#: 400x225 units times `K`. The title starts at `_TITLE_TOP`; the body box runs from
#: `_BODY_TOP` (lower by one title line for each extra line) to `_BODY_BOTTOM`, above the
#: foot rule.
_K = deck_type.K
_TITLE_TOP = 24 * _K
_BODY_TOP = deck_type.BODY_TOP * _K
_BODY_BOTTOM = deck_type.BODY_BOTTOM * _K
_FOOT_RULE = 32 * _K


def _u(name: str) -> float:
    """A size from the shared type scale, in points."""
    return deck_type.TYPE[name]


#: Default width of a picture sharing the slide with text.
_PICTURE_SPAN = 300.0


def _has_words(data: dict) -> bool:
    """Whether the slide says anything besides its title and picture."""
    return any(data.get(key) for key in ("bullets", "body", "rows", "metrics", "chart", *_PAIRED))


def _picture_span(data: dict) -> float:
    """Width of a picture sharing a slide with words, in export points."""
    return {"small": 230.0, "medium": _PICTURE_SPAN, "large": 390.0}.get(
        str((data.get("image") or {}).get("size") or "medium"), _PICTURE_SPAN
    )


_EMU_PER_PT = 12700

#: (latin, East Asian) faces keyed by `design.FONTS`; see `_font`.
_FACES = {
    "gothic": ("Segoe UI", "맑은 고딕"),
    "serif": ("Georgia", "바탕"),
}

_INK = RGBColor(0x1A, 0x1A, 0x1A)
_MUTED = RGBColor(0x66, 0x66, 0x66)

#: Dark-template ground and neutrals; not design tokens.
_DARK_BG = RGBColor(0x0E, 0x11, 0x16)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK_INK = RGBColor(0xF5, 0xF6, 0xF7)
_DARK_MUTED = RGBColor(0x9A, 0xA0, 0xA6)

#: PDF fallbacks when no design system is attached (not exactly `#1a1a1a`).
_PDF_INK = (0.1, 0.1, 0.1)
_PDF_MUTED = (0.4, 0.4, 0.4)


def _rgb(value: str | None) -> RGBColor:
    """`#rrggbb` or `#rgb` → RGBColor; the default accent for anything else."""
    text = (value or "").strip().lstrip("#")
    if len(text) == 3 and re.fullmatch(r"[0-9a-fA-F]{3}", text):
        text = "".join(c * 2 for c in text)
    if len(text) == 6 and re.fullmatch(r"[0-9a-fA-F]{6}", text):
        return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    return RGBColor(0x5B, 0x5B, 0xD6)


def _mix(
    colour: RGBColor, percent: float, *, onto: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)
) -> RGBColor:
    """`percent`% of `colour` over `onto`, matching the preview's CSS `color-mix`."""
    weight = max(0.0, min(1.0, percent / 100))
    return RGBColor(*(round(colour[i] * weight + onto[i] * (1 - weight)) for i in range(3)))


def _block(slide, *, left: float, top: float, width: float, height: float, colour: RGBColor):
    """A filled rectangle with no line and no shadow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(left * _EMU_PER_PT)),
        Emu(int(top * _EMU_PER_PT)),
        Emu(int(width * _EMU_PER_PT)),
        Emu(int(height * _EMU_PER_PT)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _hex_floats(value: str | None) -> tuple[float, float, float]:
    colour = _rgb(value)
    return colour[0] / 255, colour[1] / 255, colour[2] / 255


@dataclass(frozen=True)
class Look:
    """A visual style, matching the panel's `LOOKS`: ground, neutrals, cover, ornament, card
    style, body line height.
    """

    bg: str
    ink: str
    muted: str
    faint: str
    hair: str
    #: How much accent goes into a tint, onto the ground.
    tint: int
    card: str  # filled | outlined
    radius: float
    badge: str  # square | circle
    cover: str  # gradient | wash | glow | split | paper | brackets
    #: top-band | left-bar | corner-circle | bottom-rule | gutter | frame | bottom-band
    ornament: str
    #: Whether the neutrals above replace the design system's.
    own_neutrals: bool = False
    #: Body line height; the panel's `leading`.
    leading: float = 1.6


_LOOKS: dict[str, Look] = {
    "editorial": Look(
        bg="#ffffff",
        ink="#1a1a1a",
        muted="#666666",
        faint="#8a8a8a",
        hair="#e6e6e6",
        tint=7,
        card="filled",
        radius=0,
        badge="square",
        cover="gradient",
        ornament="top-band",
        own_neutrals=False,
    ),
    "poster": Look(
        bg="#f7f3ed",
        ink="#1a1a1a",
        muted="#666666",
        faint="#8a8a8a",
        hair="#e2ddd4",
        tint=9,
        card="filled",
        radius=0,
        badge="square",
        cover="gradient",
        ornament="left-bar",
        own_neutrals=False,
        leading=1.55,
    ),
    "minimal": Look(
        bg="#ffffff",
        ink="#1a1a1a",
        muted="#666666",
        faint="#8a8a8a",
        hair="#ececec",
        tint=5,
        card="outlined",
        radius=0,
        badge="square",
        cover="wash",
        ornament="corner-circle",
        own_neutrals=False,
        leading=1.7,
    ),
    "dark": Look(
        bg="#0f172a",
        ink="#f1f5f9",
        muted="#a3b1c6",
        faint="#64748b",
        hair="#273449",
        tint=22,
        card="filled",
        radius=6,
        badge="circle",
        cover="glow",
        ornament="bottom-rule",
        own_neutrals=True,
    ),
    "split": Look(
        bg="#ffffff",
        ink="#111827",
        muted="#5b6472",
        faint="#9aa3b2",
        hair="#e5e7eb",
        tint=6,
        card="outlined",
        radius=0,
        badge="square",
        cover="split",
        ornament="gutter",
        own_neutrals=True,
    ),
    "warm": Look(
        bg="#f6f1e8",
        ink="#3f3328",
        muted="#7a6a5a",
        faint="#a8998a",
        hair="#e2d8c8",
        tint=12,
        card="filled",
        radius=10,
        badge="circle",
        cover="paper",
        ornament="bottom-band",
        own_neutrals=True,
        leading=1.7,
    ),
    "mono": Look(
        bg="#ffffff",
        ink="#111111",
        muted="#555555",
        faint="#8a8a8a",
        hair="#111111",
        tint=0,
        card="outlined",
        radius=0,
        badge="square",
        cover="brackets",
        ornament="frame",
        own_neutrals=True,
    ),
    "pastel": Look(
        bg="#f3f0fa",
        ink="#2b2540",
        muted="#6b6480",
        faint="#9a93ad",
        hair="#e3ddf0",
        tint=14,
        card="filled",
        radius=12,
        badge="circle",
        cover="wash",
        ornament="corner-circle",
        own_neutrals=True,
        leading=1.65,
    ),
    "forest": Look(
        bg="#f1f5f0",
        ink="#1f2d22",
        muted="#5c6b5e",
        faint="#8e9a90",
        hair="#d9e2da",
        tint=10,
        card="filled",
        radius=6,
        badge="square",
        cover="gradient",
        ornament="left-bar",
        own_neutrals=True,
    ),
    "slate": Look(
        bg="#eef1f5",
        ink="#1c2431",
        muted="#55617a",
        faint="#8b96ab",
        hair="#d5dbe6",
        tint=8,
        card="outlined",
        radius=2,
        badge="square",
        cover="split",
        ornament="bottom-rule",
        own_neutrals=True,
    ),
    "paper": Look(
        bg="#fbfaf6",
        ink="#2a2622",
        muted="#6b655c",
        faint="#9b948a",
        hair="#e6e1d6",
        tint=6,
        card="outlined",
        radius=0,
        badge="square",
        cover="paper",
        ornament="frame",
        own_neutrals=True,
        leading=1.7,
    ),
}


def _look_of(visual_style: str) -> Look:
    return _LOOKS.get(visual_style) or _LOOKS["editorial"]


def _shape(slide, kind, *, left: float, top: float, width: float, height: float):
    return slide.shapes.add_shape(
        kind,
        Emu(int(left * _EMU_PER_PT)),
        Emu(int(top * _EMU_PER_PT)),
        Emu(int(width * _EMU_PER_PT)),
        Emu(int(height * _EMU_PER_PT)),
    )


def _box(
    slide,
    look: Look,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    line: RGBColor,
):
    """A card, band or metric box, filled or outlined and rounded as the look says."""
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if look.radius else MSO_SHAPE.RECTANGLE
    shape = _shape(slide, kind, left=left, top=top, width=width, height=height)
    if look.radius:
        # The adjustment is a fraction of the shorter side.
        shape.adjustments[0] = min(0.5, (look.radius * 2.4) / max(1.0, min(width, height)))
    shape.fill.solid()
    if look.card == "outlined":
        shape.fill.fore_color.rgb = _rgb(look.bg)
        shape.line.color.rgb = line
        shape.line.width = Emu(int(0.75 * _EMU_PER_PT))
    else:
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _badge(slide, look: Look, *, left: float, top: float, side: float, colour: RGBColor):
    """The numbered mark on a step or a tile: a square, or a disc."""
    kind = MSO_SHAPE.OVAL if look.badge == "circle" else MSO_SHAPE.RECTANGLE
    shape = _shape(slide, kind, left=left, top=top, width=side, height=side)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _mix_floats(
    colour: tuple[float, float, float],
    percent: float,
    *,
    onto: tuple[float, float, float] = (1.0, 1.0, 1.0),
) -> tuple[float, float, float]:
    """`_mix` for reportlab, which works in 0–1 floats rather than bytes."""
    weight = max(0.0, min(1.0, percent / 100))
    return tuple(colour[i] * weight + onto[i] * (1 - weight) for i in range(3))  # type: ignore[return-value]


def _table_size(rows: int) -> float:
    """The preview's fitted table cell size in its 225-unit drawing; both exporters scale by the
    same ratio.
    """
    return deck_type.table_size(rows)


def _font(
    run,
    *,
    size: int,
    bold: bool = False,
    colour: RGBColor = _INK,
    faces: tuple[str, str] = _FACES["gothic"],
) -> None:
    """Sets a run's font; `font.name` writes only `a:latin`, Hangul is drawn from `a:ea`."""
    latin, east_asian = faces
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = latin
    properties = run.font._rPr
    for tag in ("a:ea", "a:cs"):
        element = properties.makeelement(qn(tag), {"typeface": east_asian})
        properties.append(element)


def _pptx_pairs(
    slide,
    pairs: list[tuple[str, str]],
    *,
    layout: str,
    accent: RGBColor,
    tint: RGBColor,
    muted: RGBColor,
    width: float,
    left: float = 72.0,
    paint,
    look: Look | None = None,
    hair: RGBColor | None = None,
    top: float = _BODY_TOP,
    room: float = _BODY_BOTTOM - _BODY_TOP,
    scale: float = 1.0,
    measure: str | None = None,
    compact: bool = False,
) -> None:
    """Draws a paired layout: bands, tiles, steps, cards, or timeline.

    `compact`: the strip under a figure band — cards at 14/12 pt in the room that is left.

    `scale` is the slide's text scale, for exact line spacing; `measure` the `.pdf` face
    that tells how many lines a card name takes.
    """
    look = look or _LOOKS["editorial"]
    hair = hair or RGBColor(0xE6, 0xE6, 0xE6)
    if layout == "bands":
        label = 96.0
        height = min(72.0, (room - 10 * (len(pairs) - 1)) / max(len(pairs), 1))
        band = max(_u("bandMin"), min(_u("bandMax"), height / 3))
        for index, (name, text) in enumerate(pairs):
            y = top + index * (height + 10)
            _box(
                slide, look, left=left, top=y, width=label, height=height, fill=accent, line=accent
            )
            if look.card == "outlined":
                # The name chip stays filled.
                _block(slide, left=left, top=y, width=label, height=height, colour=accent)
            _box(
                slide,
                look,
                left=left + label + 8,
                top=y,
                width=width - label - 8,
                height=height,
                fill=tint,
                line=hair,
            )
            box = _textbox(slide, left=left, top=y + height / 2 - 12, width=label, height=24)
            box.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = box.paragraphs[0].add_run()
            run.text = name
            paint(run, size=band, bold=True, colour=_WHITE)
            body = _textbox(
                slide,
                left=left + label + 24,
                top=y + 10,
                width=width - label - 40,
                height=height - 20,
            )
            body.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = body.paragraphs[0].add_run()
            run.text = text
            paint(run, size=band)
        return

    if layout == "tiles":
        span = (width - 16 * (len(pairs) - 1)) / max(len(pairs), 1)
        side = min(span, 96.0)
        for index, (mark, name) in enumerate(pairs):
            item_left = left + index * (span + 16)
            _badge(slide, look, left=item_left, top=top + 20, side=side, colour=accent)
            box = _textbox(
                slide, left=item_left, top=top + 20 + side / 2 - 26, width=side, height=52
            )
            box.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = box.paragraphs[0].add_run()
            run.text = mark
            paint(run, size=_u("tileMark"), bold=True, colour=_WHITE)
            under = _textbox(
                slide, left=item_left - 8, top=top + 32 + side, width=side + 16, height=44
            )
            under.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = under.paragraphs[0].add_run()
            run.text = name
            paint(run, size=_u("tileName"), colour=muted)
        return

    if layout == "steps":
        gap = 18.0
        span = (width - gap * (len(pairs) - 1)) / max(len(pairs), 1)
        side = 44.0
        if len(pairs) > 1:
            # Rule from the first badge's centre to the last's.
            _block(
                slide,
                left=left + side / 2,
                top=top + 20 + side / 2 - 1,
                width=width - span,
                height=2,
                colour=tint,
            )
        for index, (name, text) in enumerate(pairs):
            item_left = left + index * (span + gap)
            _badge(slide, look, left=item_left, top=top + 20, side=side, colour=accent)
            box = _textbox(slide, left=item_left, top=top + 20 + 8, width=side, height=30)
            box.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = box.paragraphs[0].add_run()
            run.text = f"{index + 1:02d}"
            paint(run, size=_u("stepBadge"), bold=True, colour=_WHITE)
            title = _textbox(slide, left=item_left, top=top + 20 + side + 12, width=span, height=36)
            run = title.paragraphs[0].add_run()
            run.text = name
            paint(run, size=_u("stepName"), bold=True)
            body = _textbox(slide, left=item_left, top=top + 20 + side + 50, width=span, height=140)
            body.paragraphs[0].line_spacing = Pt(
                _u("stepText") * scale * deck_type.LEADING["stepText"]
            )
            run = body.paragraphs[0].add_run()
            run.text = text
            paint(run, size=_u("stepText"), colour=muted)
        return

    if layout == "cards":
        gap = 18.0
        span = (width - gap * (len(pairs) - 1)) / max(len(pairs), 1)
        height = min((60 if compact else 100) * _K, room - 10)
        name_pt = 14.0 if compact else _u("cardName")
        text_pt = 12.0 if compact else _u("cardText")
        for index, (name, text) in enumerate(pairs):
            item_left = left + index * (span + gap)
            _box(
                slide,
                look,
                left=item_left,
                top=top + 10,
                width=span,
                height=height,
                fill=tint,
                line=hair,
            )
            _block(slide, left=item_left, top=top + 10, width=span, height=5, colour=accent)
            name_lines = len(_wrap(name, measure, name_pt * scale, span - 28)) if measure else 1
            name_height = name_pt * scale * 1.3 * max(1, name_lines) + 6
            title = _textbox(
                slide, left=item_left + 14, top=top + 26, width=span - 28, height=name_height
            )
            title.paragraphs[0].line_spacing = Pt(name_pt * scale * 1.3)
            run = title.paragraphs[0].add_run()
            run.text = name
            paint(run, size=name_pt, bold=True, colour=accent)
            body = _textbox(
                slide,
                left=item_left + 14,
                top=top + 26 + name_height + 4,
                width=span - 28,
                height=height - 40 - name_height,
            )
            body.paragraphs[0].line_spacing = Pt(text_pt * scale * deck_type.LEADING["cardText"])
            run = body.paragraphs[0].add_run()
            run.text = text
            paint(run, size=text_pt)
        return

    # timeline
    axis = 128.0
    step = min(56.0, room / max(len(pairs), 1))
    line = max(_u("lineMin"), min(_u("lineMax"), step / 2.3))
    _block(slide, left=left + axis, top=top, width=1.5, height=step * len(pairs), colour=tint)
    for index, (when, what) in enumerate(pairs):
        y = top + index * step
        date = _textbox(slide, left=left, top=y, width=axis - 12, height=30)
        date.paragraphs[0].alignment = PP_ALIGN.RIGHT
        run = date.paragraphs[0].add_run()
        run.text = when
        paint(run, size=line, bold=True, colour=accent)
        _block(slide, left=left + axis - 3.5, top=y + 8, width=8, height=8, colour=accent)
        body = _textbox(slide, left=left + axis + 16, top=y, width=width - axis - 16, height=step)
        body.paragraphs[0].line_spacing = Pt(line * scale * deck_type.LEADING["line"])
        run = body.paragraphs[0].add_run()
        run.text = what
        paint(run, size=line)


def _pptx_chart(
    slide,
    chart: dict,
    *,
    accent,
    muted,
    width: float,
    faces: tuple[str, str],
    left: float = 72.0,
    top: float = _BODY_TOP,
    room: float = _BODY_BOTTOM - _BODY_TOP,
) -> None:
    """A native PowerPoint chart (editable worksheet), styled by `services.charts`."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    payload = CategoryChartData()
    payload.categories = chart["categories"]
    for name, values in chart["series"]:
        payload.add_series(name or " ", values)

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS if chart["kind"] == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(int(left * _EMU_PER_PT)),
        Emu(int(top * _EMU_PER_PT)),
        Emu(int(width * _EMU_PER_PT)),
        Emu(int(room * _EMU_PER_PT)),
        payload,
    )
    charts.apply(
        frame.chart,
        kind=chart["kind"],
        unit=chart.get("unit") or "",
        accent=accent,
        muted=muted,
        faces=faces,
    )


def _cell_rule(cell, colour: RGBColor, points: float) -> None:
    """A bottom rule on one table cell, written as `a:lnB` XML.

    `a:tcPr` is schema-ordered with line elements first; inserted at the
    front, or PowerPoint offers to repair the file.
    """
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn

    properties = cell._tc.get_or_add_tcPr()
    for existing in properties.findall(qn("a:lnB")):
        properties.remove(existing)
    properties.insert(
        0,
        parse_xml(
            f'<a:lnB {nsdecls("a")} w="{int(points * 12700)}" cap="flat"'
            ' cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{colour}"/></a:solidFill>'
            "</a:lnB>"
        ),
    )


#: Placeholder mark PowerPoint reads for the outline pane and screen readers.
#: Indices are per slide: title 0, body 1.
_PH = (
    '<p:ph xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ' type="{kind}" idx="{index}"/>'
)


def _placeholder(box, kind: str, index: int) -> None:
    """Marks a drawn textbox as a placeholder of `kind`; only placeholders are outline text."""
    from pptx.oxml import parse_xml

    box._element.nvSpPr.nvPr.append(parse_xml(_PH.format(kind=kind, index=index)))


def _textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    placeholder: tuple[str, int] | None = None,
):
    box = slide.shapes.add_textbox(
        Emu(int(left * _EMU_PER_PT)),
        Emu(int(top * _EMU_PER_PT)),
        Emu(int(width * _EMU_PER_PT)),
        Emu(int(height * _EMU_PER_PT)),
    )
    if placeholder:
        _placeholder(box, *placeholder)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _outline_placeholder(slide, kind: str, index: int, text: str, colour: RGBColor) -> None:
    """A hidden 1×1 placeholder carrying the outline text; masters re-apply placeholder geometry, so
    the visible box stays a plain textbox.
    """
    frame = _textbox(slide, left=0, top=0, width=1, height=1, placeholder=(kind, index))
    frame.clear()
    for line_index, line in enumerate(text.splitlines() or [text]):
        paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        # LibreOffice ignores `cNvPr hidden` on placeholders; 1pt in the
        # ground colour keeps it invisible there.
        run.font.size = Pt(1)
        run.font.color.rgb = colour
    frame._parent.element.nvSpPr.cNvPr.set("hidden", "1")


def _columns_of(data: dict, bullets: list[str], layout: str) -> list[list[str]]:
    """Columns to lay side by side: explicit `columns` when given, else the bullets halved."""
    given = [list(c) for c in (data.get("columns") or []) if c]
    if layout == "two-column" and len(given) >= 2:
        return given
    return _split_columns(bullets) if layout == "two-column" else [bullets]


def _split_columns(bullets: list[str]) -> list[list[str]]:
    """A list in two columns, top-to-bottom then across; one column below five items."""
    if len(bullets) < 5:
        return [bullets]
    half = (len(bullets) + 1) // 2
    return [bullets[:half], bullets[half:]]


#: Logo height in the foot, and its width cap.
_LOGO_HEIGHT = 18.0
_LOGO_MAX_WIDTH = 120.0


def _logo_of(tokens: dict[str, str] | None) -> tuple[bytes, float, float] | None:
    """`(bytes, width, height)` for the design system's logo drawn to `_LOGO_HEIGHT`, or None."""
    raw = (tokens or {}).get("logo") or ""
    if not raw.startswith("data:image/"):
        return None
    try:
        blob = base64.b64decode(raw.split(",", 1)[1], validate=False)
        with PIL.Image.open(io.BytesIO(blob)) as picture:
            width, height = picture.size
    except Exception:  # noqa: BLE001 — a mark that will not decode is no mark
        log.warning("could not read the design system's logo")
        return None
    if not width or not height:
        return None
    drawn_width = min(_LOGO_MAX_WIDTH, _LOGO_HEIGHT * width / height)
    return blob, drawn_width, _LOGO_HEIGHT * min(1.0, _LOGO_MAX_WIDTH / max(drawn_width, 1e-6))


#: Paired layouts; see `deck._PAIRED`.
_PAIRED = ("bands", "tiles", "timeline", "steps", "cards")

#: Layouts drawn reversed out of the accent.
_COVERS = ("title", "section", "closing")


def _written(slides: list[dict]) -> list[dict]:
    """Slides minus those still marked `deck.UNWRITTEN`; numbering is by position."""
    return [
        slide
        for slide in slides
        if not (str(slide.get("body") or "").strip() == deck.UNWRITTEN and not _filled(slide))
    ]


def _filled(slide: dict) -> bool:
    """Whether anything but the placeholder is on this slide."""
    return any(slide.get(key) for key in ("bullets", "rows", "metrics", "chart", "image", *_PAIRED))


def _pairs_of(slide: dict, layout: str) -> list[tuple[str, str]]:
    """`[(왼쪽, 오른쪽)]` for a paired layout with both halves filled; empty otherwise."""
    if layout not in _PAIRED:
        return []
    out: list[tuple[str, str]] = []
    for item in slide.get(layout) or []:
        pair = item if isinstance(item, (list, tuple)) else ()
        if len(pair) >= 2 and str(pair[0]).strip() and str(pair[1]).strip():
            out.append((str(pair[0]).strip(), str(pair[1]).strip()))
    return out


def _chart_of(slide: dict) -> dict | None:
    """A drawable chart, or None; categories and series are cut to the shortest paired length."""
    chart = slide.get("chart")
    if not isinstance(chart, dict):
        return None
    categories = [str(c) for c in (chart.get("categories") or [])]
    series: list[tuple[str, list[float]]] = []
    for item in chart.get("series") or []:
        if not isinstance(item, dict):
            continue
        values: list[float] = []
        for raw in item.get("values") or []:
            try:
                values.append(float(raw))
            except (TypeError, ValueError):
                break
        if values:
            series.append((str(item.get("name") or ""), values))
    if not categories or not series:
        return None
    width = min(len(categories), min(len(values) for _, values in series))
    if width < 2:
        return None
    return {
        "kind": "line" if str(chart.get("kind")) == "line" else "bar",
        "unit": str(chart.get("unit") or ""),
        "categories": categories[:width],
        "series": [(name, values[:width]) for name, values in series],
    }


def _picture_of(slide: dict) -> tuple[bytes, str] | None:
    """`(bytes, caption)` from `image.data` (bytes, via `page_export`) or `image.src` (a `data:`
    URI).
    """
    raw = slide.get("image")
    if not isinstance(raw, dict):
        return None
    caption = str(raw.get("caption") or "")
    data = raw.get("data")
    if isinstance(data, bytes | bytearray) and data:
        return bytes(data), caption
    decoded = pictures.decode(str(raw.get("src") or ""))
    return (decoded[1], caption) if decoded else None


def _fit(data: bytes, *, box: tuple[float, float]) -> tuple[float, float]:
    """The size a picture takes inside `box`, in points, keeping its aspect."""
    max_width, max_height = box
    try:
        with PIL.Image.open(io.BytesIO(data)) as picture:
            width, height = picture.size
    except Exception:  # noqa: BLE001 — a picture we cannot measure gets the box
        return max_width, max_height
    if not width or not height:
        return max_width, max_height
    scale = min(max_width / width, max_height / height)
    return width * scale, height * scale


def _fill(
    data: bytes, *, box: tuple[float, float]
) -> tuple[float, float, float, float, float, float]:
    """`(width, height, left, top, right, bottom)`: the cover-fit size and PowerPoint's 0–1 crop
    fractions.
    """
    box_width, box_height = box
    try:
        with PIL.Image.open(io.BytesIO(data)) as picture:
            width, height = picture.size
    except Exception:  # noqa: BLE001 — the caller already handles bad bytes
        return box_width, box_height, 0.0, 0.0, 0.0, 0.0
    if not width or not height:
        return box_width, box_height, 0.0, 0.0, 0.0, 0.0
    scale = max(box_width / width, box_height / height)
    drawn_width, drawn_height = width * scale, height * scale
    horizontal = max(0.0, (drawn_width - box_width) / drawn_width / 2)
    vertical = max(0.0, (drawn_height - box_height) / drawn_height / 2)
    return drawn_width, drawn_height, horizontal, vertical, horizontal, vertical


def to_pptx(
    title: str,
    slides: list[dict],
    *,
    tokens: dict[str, str] | None = None,
    dark: bool = False,
    template: str = "",
) -> bytes:
    """The deck as a PowerPoint file, drawn on the blank layout.

    `tokens` is the design system copied onto the artifact; `template` is the
    서식's `.pptx`, whose master and theme the file is built on.
    """
    style = design.normalise_tokens(tokens) if tokens else None
    faces = _FACES[style["font"]] if style else _FACES["gothic"]
    ink = _rgb(style["ink"]) if style else _INK
    muted = _rgb(style["muted"]) if style else _MUTED
    if dark:
        # Design-system ink is for paper; on a dark ground the neutrals swap.
        ink, muted = _DARK_INK, _DARK_MUTED

    mark = _logo_of(style)
    footer = (style or {}).get("footer") or ""

    #: Current slide's `textScale`; a list so `paint` can read the rebinding.
    typescale = [1.0]
    #: The `.pdf` face, for measuring how many lines a title takes.
    measure = fonts.korean(style["font"] if style else "gothic")

    def paint(
        run,
        *,
        size: float,
        bold: bool = False,
        colour: RGBColor | None = None,
        fixed: bool = False,
    ) -> None:
        # Titles keep their size; everything else follows the slide's scale, never
        # under the floor.
        factor = 1.0 if fixed else typescale[0]
        _font(
            run,
            size=max(deck_type.FLOOR_PT, round(size * factor)),
            bold=bold,
            colour=colour or ink,
            faces=faces,
        )

    def paint_rich(
        paragraph,
        data: dict,
        key: str,
        text: str,
        *,
        size: float,
        bold: bool = False,
        colour: RGBColor | None = None,
        fixed: bool = False,
    ) -> None:
        html = (data.get("richText") or {}).get(key)
        scale = {1: 0.65, 2: 0.8, 3: 1.0, 4: 1.2, 5: 1.5, 6: 2.0, 7: 3.0}
        for value, inline in _inline_runs(html, text):
            run = paragraph.add_run()
            run.text = value
            run_colour = colour
            raw_colour = str(inline.get("color") or "")
            if re.fullmatch(r"#[0-9a-fA-F]{6}", raw_colour):
                run_colour = _rgb(raw_colour)
            paint(
                run,
                size=max(
                    8,
                    round(size * float(inline.get("scale") or scale.get(inline.get("size"), 1.0))),
                ),
                bold=bool(inline.get("bold", bold)),
                colour=run_colour,
                fixed=fixed,
            )
            run.font.italic = bool(inline.get("italic"))
            run.font.underline = bool(inline.get("underline"))

    # Slides are drawn as free shapes, not placed in the template's placeholders.
    presentation = Presentation(template) if template else Presentation()
    presentation.slide_width = Emu(int(_W * _EMU_PER_PT))
    presentation.slide_height = Emu(int(_H * _EMU_PER_PT))
    blank = presentation.slide_layouts[6]
    visual_style = (style or {}).get("visualStyle") or "editorial"

    for index, data in enumerate(_written(slides)):
        slide = presentation.slides.add_slide(blank)
        typescale[0] = _typescale(data)
        if dark:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _DARK_BG
        look = _look_of(visual_style)
        if look.own_neutrals and not dark:
            ink, muted = _rgb(look.ink), _rgb(look.muted)
        # The slide's own accent, else the design system's.
        accent = _rgb(data.get("accent") or (style or {}).get("accent"))
        if visual_style == "mono":
            accent = ink
        ground = _rgb(look.bg)
        tint = _mix(accent, look.tint, onto=ground) if look.tint else RGBColor(0xF2, 0xF2, 0xF2)
        hair = _rgb(look.hair)
        layout = data.get("layout") or "bullets"

        cover = layout in _COVERS
        #: Cover text on the accent (white ink) rather than the look's ground.
        on_accent = look.cover in ("gradient", "glow")
        if cover:
            fill = slide.background.fill
            if look.cover == "gradient":
                try:
                    fill.gradient()
                    fill.gradient_angle = 45.0
                    stops = fill.gradient_stops
                    stops[0].color.rgb = accent
                    stops[1].color.rgb = _mix(
                        accent, 48 if visual_style == "poster" else 62, onto=_INK
                    )
                    for extra in list(stops)[2:]:
                        extra.color.rgb = stops[1].color.rgb
                except Exception as exc:  # noqa: BLE001 — a ground, not the deck
                    log.warning("could not fill a cover with a gradient: %s", exc)
                    fill.solid()
                    fill.fore_color.rgb = accent
            elif look.cover == "wash":
                fill.solid()
                fill.fore_color.rgb = _mix(accent, 10)
            else:
                fill.solid()
                fill.fore_color.rgb = ground
                if look.cover == "glow":
                    for ring, pct in ((520, 22), (400, 38), (280, 58)):
                        disc = _shape(
                            slide,
                            MSO_SHAPE.OVAL,
                            left=_W - ring * 0.55,
                            top=_H - ring * 0.45,
                            width=ring,
                            height=ring,
                        )
                        disc.fill.solid()
                        disc.fill.fore_color.rgb = _mix(accent, pct, onto=ground)
                        disc.line.fill.background()
                        disc.shadow.inherit = False
                elif look.cover == "paper":
                    disc = _shape(
                        slide, MSO_SHAPE.OVAL, left=_W - 300, top=44, width=456, height=456
                    )
                    disc.fill.solid()
                    disc.fill.fore_color.rgb = accent
                    disc.line.fill.background()
                    disc.shadow.inherit = False
                elif look.cover == "split":
                    _block(slide, left=0, top=0, width=384, height=_H, colour=accent)
                elif look.cover == "brackets":
                    ink_line = ink
                    _block(slide, left=62, top=67, width=6, height=72, colour=ink_line)
                    _block(slide, left=62, top=67, width=53, height=6, colour=ink_line)
                    _block(slide, left=_W - 68, top=_H - 139, width=6, height=72, colour=ink_line)
                    _block(slide, left=_W - 115, top=_H - 73, width=53, height=6, colour=ink_line)
        else:
            if look.bg.lower() != "#ffffff":
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = ground
            if look.ornament == "left-bar":
                _block(slide, left=0, top=0, width=14, height=_H, colour=accent)
            elif look.ornament == "top-band":
                _block(slide, left=0, top=0, width=_W, height=14, colour=accent)
            elif look.ornament == "corner-circle":
                disc = _shape(slide, MSO_SHAPE.OVAL, left=_W - 96, top=-84, width=168, height=168)
                disc.fill.solid()
                disc.fill.fore_color.rgb = _mix(accent, 12, onto=ground)
                disc.line.fill.background()
                disc.shadow.inherit = False
            elif look.ornament == "bottom-rule":
                _block(slide, left=0, top=_H - 10, width=_W * 0.6, height=10, colour=accent)
                _block(
                    slide,
                    left=_W * 0.6,
                    top=_H - 10,
                    width=_W * 0.4,
                    height=10,
                    colour=_mix(accent, 40, onto=ground),
                )
            elif look.ornament == "gutter":
                _block(slide, left=0, top=0, width=7, height=_H, colour=accent)
                counter = _textbox(slide, left=22, top=_H - 118, width=80, height=50)
                run = counter.paragraphs[0].add_run()
                run.text = f"{index + 1:02d}"
                paint(run, size=_u("gutterNumber"), bold=True, colour=accent)
            elif look.ornament == "frame":
                frame_line = _shape(
                    slide, MSO_SHAPE.RECTANGLE, left=24, top=24, width=_W - 48, height=_H - 48
                )
                frame_line.fill.background()
                frame_line.line.color.rgb = ink
                frame_line.line.width = Emu(int(1.0 * _EMU_PER_PT))
                frame_line.shadow.inherit = False
            elif look.ornament == "bottom-band":
                _block(slide, left=0, top=_H - 53, width=_W, height=53, colour=tint)

        heading = str(data.get("title") or "")
        body = str(data.get("body") or "")
        bullets = [str(b) for b in (data.get("bullets") or []) if str(b).strip()]
        rows = [[str(cell) for cell in row] for row in (data.get("rows") or []) if row]
        metrics = [
            [str(pair[0]), str(pair[1])]
            for pair in (data.get("metrics") or [])
            if isinstance(pair, list) and len(pair) >= 2
        ]
        chart = _chart_of(data)
        picture = _picture_of(data)
        picture_span = _picture_span(data)
        pairs = _pairs_of(data, layout)
        picture_left = bool(
            picture
            and (bullets or rows or metrics or chart or body)
            and str((data.get("image") or {}).get("position") or "") == "left"
        )
        # A picture beside text narrows the text column; alone, it is centred.
        text_width = (
            _W
            - 144
            - (
                picture_span + 24
                if picture and (bullets or rows or metrics or chart or body)
                else 0
            )
        )
        text_left = 72 + (picture_span + 24 if picture_left else 0)

        #: Past the accent block on a split cover.
        cover_left = 72 + 336 if look.cover == "split" else 72
        cover_size = _u(
            "coverPoster"
            if visual_style == "poster"
            else "coverMono"
            if visual_style == "mono"
            else "cover"
        )
        #: Body box for this slide; a long title lowers its top.
        body_top, room = _BODY_TOP, _BODY_BOTTOM - _BODY_TOP
        if layout == "closing":
            cover_ink = _WHITE if on_accent else ink
            cover_muted = _mix(_WHITE, 80, onto=accent) if on_accent else muted
            if look.cover == "split":
                counter = _textbox(slide, left=60, top=_H - 130, width=200, height=80)
                run = counter.paragraphs[0].add_run()
                run.text = "END"
                paint(run, size=_u("splitNumber"), bold=True, colour=_mix(_WHITE, 35, onto=accent))
            if look.cover != "brackets":
                _block(
                    slide,
                    left=cover_left,
                    top=120,
                    width=106,
                    height=7,
                    colour=_WHITE if on_accent else accent,
                )
            closing_size = _u("closing") * typescale[0]
            closing_lines = len(
                _wrap(heading or "마무리", measure, closing_size, _W - 72 - cover_left)
            )
            frame = _textbox(
                slide,
                left=cover_left,
                top=140,
                width=_W - 72 - cover_left,
                height=closing_size * 1.2 * max(1, closing_lines) + 8,
                placeholder=("title", 0),
            )
            frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            paint_rich(
                frame.paragraphs[0],
                data,
                "title",
                heading or "마무리",
                size=_u("closing"),
                bold=True,
                colour=cover_ink,
            )
            if bullets:
                listing = _textbox(
                    slide,
                    left=cover_left,
                    top=140 + closing_size * 1.2 * max(1, closing_lines) + 20,
                    width=_W - 72 - cover_left,
                    height=200,
                    placeholder=("body", 1),
                )
                for position, text in enumerate(bullets[:3]):
                    paragraph = listing.paragraphs[0] if position == 0 else listing.add_paragraph()
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.space_after = Pt(10)
                    marker = paragraph.add_run()
                    marker.text = "— "
                    paint(marker, size=_u("closingBullets"), bold=True, colour=cover_muted)
                    paint_rich(
                        paragraph,
                        data,
                        f"bullets.{position}",
                        text,
                        size=_u("closingBullets"),
                        colour=cover_ink,
                    )
            if body:
                foot = _textbox(slide, left=cover_left, top=_H - 120, width=_W - 144, height=50)
                foot.paragraphs[0].alignment = PP_ALIGN.LEFT
                paint_rich(
                    foot.paragraphs[0],
                    data,
                    "body",
                    body,
                    size=_u("closingBody"),
                    bold=True,
                    colour=cover_ink,
                )
        elif cover:
            cover_ink = _WHITE if on_accent else ink
            cover_muted = _mix(_WHITE, 80, onto=accent) if on_accent else muted
            if look.cover == "split":
                counter = _textbox(slide, left=60, top=_H - 130, width=200, height=80)
                run = counter.paragraphs[0].add_run()
                number_text = str(data.get("number") or "").replace(".", "") or "01"
                run.text = "END" if layout == "closing" else number_text
                paint(run, size=_u("splitNumber"), bold=True, colour=_mix(_WHITE, 35, onto=accent))
            if (
                look.cover != "split"
                and layout == "section"
                and (number := str(data.get("number") or ""))
            ):
                counter = _textbox(slide, left=72, top=150, width=200, height=40)
                run = counter.paragraphs[0].add_run()
                run.text = number
                paint(
                    run,
                    size=_u("sectionNumber"),
                    bold=True,
                    colour=_mix(_WHITE, 70, onto=accent) if on_accent else accent,
                )
            elif look.cover != "brackets":
                _block(
                    slide,
                    left=cover_left + 10,
                    top=186,
                    width=106,
                    height=7,
                    colour=_WHITE if on_accent else accent,
                )
            frame = _textbox(
                slide,
                left=cover_left,
                top=210,
                width=(_W - 72 - cover_left) - (300 if look.cover == "paper" else 0),
                height=180,
            )
            frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            paint_rich(
                frame.paragraphs[0],
                data,
                "title",
                heading or title,
                size=cover_size,
                bold=True,
                colour=cover_ink,
            )
            if body:
                paragraph = frame.add_paragraph()
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_before = Pt(14)
                paint_rich(paragraph, data, "body", body, size=_u("coverBody"), colour=cover_muted)
            _outline_placeholder(
                slide,
                "ctrTitle",
                0,
                "\n".join(part for part in (heading or title, body) if part),
                accent if on_accent else ground,
            )
        elif layout == "statement":
            _block(slide, left=(_W - 62) / 2, top=176, width=62, height=5, colour=accent)
            frame = _textbox(
                slide, left=90, top=196, width=_W - 180, height=160, placeholder=("title", 0)
            )
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            paint_rich(
                frame.paragraphs[0],
                data,
                "title",
                heading,
                size=_u("statement"),
                bold=True,
                colour=accent,
            )
            if body:
                under = _textbox(
                    slide, left=120, top=360, width=_W - 240, height=80, placeholder=("body", 1)
                )
                under.paragraphs[0].alignment = PP_ALIGN.CENTER
                paint_rich(
                    under.paragraphs[0], data, "body", body, size=_u("statementBody"), colour=muted
                )
        elif layout == "quote":
            frame = _textbox(
                slide,
                left=90,
                top=170,
                width=_W - 180,
                height=200,
                placeholder=("title", 0),
            )
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = f"“{body}”" if body else heading
            paint(run, size=_u("quote"), bold=True, colour=accent)
            if body and heading:
                caption = frame.add_paragraph()
                caption.space_before = Pt(16)
                run = caption.add_run()
                run.text = heading
                paint(run, size=_u("quoteBy"), colour=muted)
        else:
            title_size = deck_type.title_pt(heading, text_width / _K)
            title_line = title_size * deck_type.LEADING["title"]
            title_lines = max(1, len(_wrap(heading, measure, title_size, text_width)))
            frame = _textbox(
                slide,
                left=text_left,
                top=_TITLE_TOP,
                width=text_width,
                height=title_line * title_lines + 6,
                placeholder=("title", 0),
            )
            frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            frame.paragraphs[0].line_spacing = Pt(title_line)
            paint_rich(
                frame.paragraphs[0], data, "title", heading, size=title_size, bold=True, fixed=True
            )
            # The tab under the title; the body box starts below it and loses one line for
            # each extra title line.
            tab_top = _TITLE_TOP + title_line * title_lines + 6 * _K
            _block(slide, left=text_left, top=tab_top, width=62, height=5, colour=accent)
            body_top = _BODY_TOP + title_line * (title_lines - 1)
            room = _BODY_BOTTOM - body_top

            # A figure the deck drew for itself is a band above the words, full width;
            # the words beneath are set compact. `picture` is spent here.
            compact = False
            if picture and (data.get("image") or {}).get("diagram") and _has_words(data):
                image_bytes, _caption = picture
                text_width, text_left = _W - 144, 72.0
                band_width, band_height = _fit(image_bytes, box=(text_width, room * 0.56))
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(image_bytes),
                        Emu(int((text_left + (text_width - band_width) / 2) * _EMU_PER_PT)),
                        Emu(int(body_top * _EMU_PER_PT)),
                        Emu(int(band_width * _EMU_PER_PT)),
                        Emu(int(band_height * _EMU_PER_PT)),
                    )
                except Exception as exc:  # noqa: BLE001 — a bad picture is not a failed export
                    log.warning("could not place a figure band into the deck: %s", exc)
                body_top += band_height + 8
                room -= band_height + 8
                picture = None
                compact = True

            if pairs:
                _pptx_pairs(
                    slide,
                    pairs,
                    layout=layout,
                    accent=accent,
                    tint=tint,
                    muted=muted,
                    width=text_width,
                    left=text_left,
                    paint=paint,
                    look=look,
                    hair=hair,
                    top=body_top,
                    room=room,
                    scale=typescale[0],
                    measure=measure,
                    compact=compact,
                )
            elif chart:
                _pptx_chart(
                    slide,
                    chart,
                    accent=accent,
                    muted=muted,
                    width=text_width,
                    faces=faces,
                    left=text_left,
                    top=body_top,
                    room=room,
                )
            elif metrics and layout == "big-number":
                figure, label = metrics[0]
                box = _textbox(
                    slide,
                    left=text_left,
                    top=body_top,
                    width=text_width,
                    height=140,
                    placeholder=("body", 1),
                )
                run = box.paragraphs[0].add_run()
                run.text = figure
                paint(run, size=_u("bigNumber"), bold=True, colour=accent)
                run = box.paragraphs[0].add_run()
                run.text = f"  {label}"
                paint(run, size=_u("bigNumberLabel"), colour=muted)
                if body:
                    under = _textbox(
                        slide, left=text_left, top=body_top + 150, width=text_width, height=80
                    )
                    paint_rich(under.paragraphs[0], data, "body", body, size=_u("bigNumberBody"))
            elif bullets and layout == "agenda":
                # Two columns above four entries.
                columns = _split_columns(bullets) if len(bullets) > 4 else [bullets]
                span = (text_width - 24 * (len(columns) - 1)) / len(columns)
                per = max(len(column) for column in columns)
                step = min(66.0, room / max(per, 1))
                number = 0
                for column_index, column in enumerate(columns):
                    left = text_left + column_index * (span + 24)
                    for position, text in enumerate(column):
                        number += 1
                        y = body_top + position * step
                        counter = _textbox(slide, left=left, top=y, width=64, height=step)
                        run = counter.paragraphs[0].add_run()
                        run.text = f"{number:02d}"
                        paint(run, size=_u("agendaNumber"), bold=True, colour=accent)
                        name = _textbox(
                            slide,
                            left=left + 64,
                            top=y + 4,
                            width=span - 64,
                            height=step,
                            placeholder=("body", 1) if number == 1 else None,
                        )
                        run = name.paragraphs[0].add_run()
                        run.text = text
                        paint(run, size=_u("agenda"))
                        _block(
                            slide, left=left, top=y + step - 6, width=span, height=0.75, colour=hair
                        )
            elif metrics:
                span = (text_width - 24 * (len(metrics) - 1)) / len(metrics)
                for position, (figure, label) in enumerate(metrics):
                    left = text_left + position * (span + 24)
                    _box(
                        slide,
                        look,
                        left=left,
                        top=body_top + 8,
                        width=span,
                        height=150,
                        fill=tint,
                        line=hair,
                    )
                    _block(slide, left=left, top=body_top + 8, width=span, height=5, colour=accent)
                    box = _textbox(
                        slide,
                        left=left,
                        top=body_top + 20,
                        width=span,
                        height=130,
                    )
                    run = box.paragraphs[0].add_run()
                    run.text = figure
                    paint(run, size=_u("metric"), bold=True, colour=accent)
                    under = box.add_paragraph()
                    under.space_before = Pt(6)
                    run = under.add_run()
                    run.text = label
                    paint(run, size=_u("metricLabel"), colour=muted)
            elif rows:
                row_height = deck_type.table_row_height(len(rows)) * _K
                shape = slide.shapes.add_table(
                    len(rows),
                    max(len(row) for row in rows),
                    Emu(int(text_left * _EMU_PER_PT)),
                    Emu(int(body_top * _EMU_PER_PT)),
                    Emu(int(text_width * _EMU_PER_PT)),
                    Emu(int(row_height * len(rows) * _EMU_PER_PT)),
                )
                table = shape.table
                # Columns by their widest cell, rows by the shared row height; the panel
                # draws the same shares.
                for column_index, share in enumerate(deck_type.column_shares(rows)):
                    table.columns[column_index].width = Emu(int(text_width * share * _EMU_PER_PT))
                for table_row in table.rows:
                    table_row.height = Emu(int(row_height * _EMU_PER_PT))
                cell_size = deck_type.table_size(len(rows)) * _K
                cell_pad = deck_type.table_pad(len(rows)) * _K
                # Off, or PowerPoint applies its theme's banded table style.
                table.first_row = False
                table.horz_banding = False
                for r, row in enumerate(rows):
                    for c, text in enumerate(row):
                        if c >= len(table.columns):
                            continue
                        cell = table.cell(r, c)
                        if r == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = accent
                        elif r % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = tint
                        else:
                            cell.fill.background()
                        if r == 0:
                            pass
                        elif r < len(rows) - 1:
                            _cell_rule(cell, muted, 0.5)
                        cell.text = ""
                        cell.margin_left = cell.margin_right = Emu(int(9 * _K * _EMU_PER_PT))
                        cell.margin_top = cell.margin_bottom = Emu(int(cell_pad * _EMU_PER_PT))
                        cell.text_frame.paragraphs[0].line_spacing = Pt(
                            cell_size * typescale[0] * deck_type.LEADING["table"]
                        )
                        run = cell.text_frame.paragraphs[0].add_run()
                        run.text = text
                        paint(
                            run,
                            size=cell_size,
                            bold=r == 0 or c == 0,
                            colour=_WHITE if r == 0 else ink,
                        )
            elif bullets:
                # PowerPoint has no column flow, so columns are separate boxes.
                columns = _columns_of(data, bullets, layout)
                span = (text_width - (24 * (len(columns) - 1))) / len(columns)
                size = _u("bodyNarrow") if len(columns) > 1 else _u("body")
                drawn = max(deck_type.FLOOR_PT, size * typescale[0])
                line_pt = drawn * look.leading
                gap_pt = drawn * deck_type.BULLET_GAP
                bullet_at = 0
                for column_index, column in enumerate(columns):
                    listing = _textbox(
                        slide,
                        left=text_left + column_index * (span + 24),
                        top=body_top,
                        width=span,
                        height=room,
                        # Only the first column is the body placeholder.
                        placeholder=("body", 1) if column_index == 0 else None,
                    )
                    for position, text in enumerate(column):
                        paragraph = (
                            listing.paragraphs[0] if position == 0 else listing.add_paragraph()
                        )
                        paragraph.line_spacing = Pt(line_pt)
                        paragraph.space_after = Pt(gap_pt)
                        marker = paragraph.add_run()
                        marker.text = "• "
                        paint(marker, size=size, bold=True, colour=accent)
                        paint_rich(paragraph, data, f"bullets.{bullet_at}", text, size=size)
                        bullet_at += 1
            elif body:
                paragraph_frame = _textbox(
                    slide,
                    left=text_left,
                    top=body_top,
                    width=text_width,
                    height=room,
                    placeholder=("body", 1),
                )
                drawn = max(deck_type.FLOOR_PT, _u("paragraph") * typescale[0])
                paragraph_frame.paragraphs[0].line_spacing = Pt(drawn * look.leading)
                paint_rich(
                    paragraph_frame.paragraphs[0],
                    data,
                    "body",
                    body,
                    size=_u("paragraph"),
                    colour=muted,
                )

        if picture:
            image_bytes, image_caption = picture
            alone = not (bullets or rows or metrics or chart or body)
            box = (_W - 260, room) if alone else (picture_span, room)
            fill = str((data.get("image") or {}).get("fit") or "") == "cover"
            if fill:
                _, _, crop_left, crop_top, crop_right, crop_bottom = _fill(image_bytes, box=box)
                width, height = box
            else:
                width, height = _fit(image_bytes, box=box)
                crop_left = crop_top = crop_right = crop_bottom = 0.0
            left = (72 if picture_left else 72 + text_width + 24) if not alone else (_W - width) / 2
            top = body_top + max(0.0, (room - height) / 2)
            try:
                shape = slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    Emu(int(left * _EMU_PER_PT)),
                    Emu(int(top * _EMU_PER_PT)),
                    Emu(int(width * _EMU_PER_PT)),
                    Emu(int(height * _EMU_PER_PT)),
                )
                if fill:
                    shape.crop_left = crop_left
                    shape.crop_top = crop_top
                    shape.crop_right = crop_right
                    shape.crop_bottom = crop_bottom
            except Exception as exc:  # noqa: BLE001 — one bad picture, not a failed export
                log.warning("could not place a picture in the pptx: %s", exc)
            else:
                if image_caption:
                    frame = _textbox(
                        slide, left=left, top=top + height + 6, width=max(width, 120), height=24
                    )
                    run = frame.paragraphs[0].add_run()
                    run.text = image_caption
                    paint(run, size=_u("caption"), colour=muted)

        # Foot: logo, deck title and footer on the left, slide number on the right.
        if not cover:
            _block(slide, left=72, top=_H - 58, width=_W - 144, height=0.75, colour=hair)
            edge = 72.0
            if mark:
                blob, mark_width, mark_height = mark
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(blob),
                        Emu(int(edge * _EMU_PER_PT)),
                        Emu(int((_H - 60) * _EMU_PER_PT)),
                        Emu(int(mark_width * _EMU_PER_PT)),
                        Emu(int(mark_height * _EMU_PER_PT)),
                    )
                except Exception as exc:  # noqa: BLE001 — a mark, not the deck
                    log.warning("could not place the logo on a pptx slide: %s", exc)
                else:
                    edge += mark_width + 10
            _block(slide, left=72, top=_H - _FOOT_RULE, width=_W - 144, height=0.75, colour=hair)
            name = _textbox(slide, left=edge, top=_H - 62, width=_W - 320 - edge, height=30)
            run = name.paragraphs[0].add_run()
            run.text = title
            paint(run, size=_u("footer"), colour=muted)
            if footer:
                who = _textbox(slide, left=_W - 330, top=_H - 62, width=210, height=30)
                who.paragraphs[0].alignment = PP_ALIGN.RIGHT
                run = who.paragraphs[0].add_run()
                run.text = footer
                paint(run, size=_u("footer"), colour=muted)
            chip = _block(slide, left=_W - 108, top=_H - 60, width=36, height=36, colour=accent)
            frame = chip.text_frame
            frame.word_wrap = False
            frame.margin_left = frame.margin_right = 0
            frame.margin_top = frame.margin_bottom = 0
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = frame.paragraphs[0].add_run()
            run.text = str(index + 1)
            paint(run, size=_u("pageNumber"), bold=True, colour=_WHITE)

        notes = str(data.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _pdf_box(
    pdf,
    look: Look,
    *,
    left: float,
    bottom: float,
    width: float,
    height: float,
    fill,
    line,
    bg,
) -> None:
    """The `.pdf` twin of `_box`."""
    if look.card == "outlined":
        pdf.setFillColorRGB(*bg)
        pdf.setStrokeColorRGB(*line)
        pdf.setLineWidth(0.75)
        if look.radius:
            pdf.roundRect(left, bottom, width, height, look.radius * 1.2, stroke=1, fill=1)
        else:
            pdf.rect(left, bottom, width, height, stroke=1, fill=1)
        return
    pdf.setFillColorRGB(*fill)
    if look.radius:
        pdf.roundRect(left, bottom, width, height, look.radius * 1.2, stroke=0, fill=1)
    else:
        pdf.rect(left, bottom, width, height, stroke=0, fill=1)


def _pdf_badge(pdf, look: Look, *, left: float, bottom: float, side: float, colour) -> None:
    pdf.setFillColorRGB(*colour)
    if look.badge == "circle":
        pdf.circle(left + side / 2, bottom + side / 2, side / 2, stroke=0, fill=1)
    else:
        pdf.rect(left, bottom, side, side, stroke=0, fill=1)


def _pdf_pairs(
    pdf,
    pairs: list[tuple[str, str]],
    *,
    layout: str,
    accent,
    tint,
    muted,
    ink,
    top: float,
    width: float,
    font: str,
    scale: float,
    left: float = 72.0,
    look: Look | None = None,
    hair=None,
    bg=None,
    room: float | None = None,
    bold: str | None = None,
    compact: bool = False,
) -> None:
    """The `.pdf` twin of `_pptx_pairs`: same geometry, explicit wrapping."""
    bold = bold or font
    room = room if room is not None else top - (_H - _BODY_BOTTOM)
    look = look or _LOOKS["editorial"]
    hair = hair or (0.902, 0.902, 0.902)
    bg = bg or (1.0, 1.0, 1.0)

    def S(n: float) -> float:
        return n * scale

    if layout == "bands":
        label = 96.0
        height = min(72.0, (room - 10 * (len(pairs) - 1)) / max(len(pairs), 1))
        band = max(_u("bandMin"), min(_u("bandMax"), height / 3))
        lead = band * deck_type.LEADING["band"]
        for index, (name, text) in enumerate(pairs):
            bottom = top - height - index * (height + 10)
            pdf.setFillColorRGB(*accent)
            pdf.rect(left, bottom, label, height, stroke=0, fill=1)
            _pdf_box(
                pdf,
                look,
                left=left + label + 8,
                bottom=bottom,
                width=width - label - 8,
                height=height,
                fill=tint,
                line=hair,
                bg=bg,
            )
            pdf.setFillColorRGB(1, 1, 1)
            pdf.setFont(bold, S(band))
            pdf.drawCentredString(left + label / 2, bottom + height / 2 - S(band) * 0.35, name)
            pdf.setFillColorRGB(*ink)
            pdf.setFont(font, S(band))
            lines = _wrap(text, font, S(band), width - label - 40)[:3]
            line_top = bottom + height / 2 + S(lead) * (len(lines) - 1) / 2 - S(band) * 0.35
            for offset, line in enumerate(lines):
                pdf.drawString(left + label + 24, line_top - offset * S(lead), line)
        return

    if layout == "tiles":
        span = (width - 16 * (len(pairs) - 1)) / max(len(pairs), 1)
        side = min(span, 96.0)
        for index, (mark, name) in enumerate(pairs):
            item_left = left + index * (span + 16)
            _pdf_badge(pdf, look, left=item_left, bottom=top - side - 20, side=side, colour=accent)
            pdf.setFillColorRGB(1, 1, 1)
            pdf.setFont(bold, S(_u("tileMark")))
            pdf.drawCentredString(
                item_left + side / 2, top - side / 2 - 20 - S(_u("tileMark")) * 0.35, mark
            )
            pdf.setFillColorRGB(*muted)
            pdf.setFont(font, S(_u("tileName")))
            for offset, line in enumerate(_wrap(name, font, S(_u("tileName")), side + 16)[:2]):
                pdf.drawCentredString(
                    item_left + side / 2,
                    top - side - 20 - S(_u("tileName")) * 1.5 - offset * S(_u("tileName")) * 1.5,
                    line,
                )
        return

    if layout == "steps":
        gap = 18.0
        span = (width - gap * (len(pairs) - 1)) / max(len(pairs), 1)
        side = 44.0
        square_top = top - 20
        if len(pairs) > 1:
            pdf.setFillColorRGB(*tint)
            pdf.rect(left + side / 2, square_top - side / 2 - 1, width - span, 2, stroke=0, fill=1)
        for index, (name, text) in enumerate(pairs):
            item_left = left + index * (span + gap)
            _pdf_badge(
                pdf, look, left=item_left, bottom=square_top - side, side=side, colour=accent
            )
            pdf.setFillColorRGB(1, 1, 1)
            pdf.setFont(bold, S(_u("stepBadge")))
            pdf.drawCentredString(
                item_left + side / 2,
                square_top - side / 2 - S(_u("stepBadge")) * 0.35,
                f"{index + 1:02d}",
            )
            pdf.setFillColorRGB(*ink)
            pdf.setFont(bold, S(_u("stepName")))
            name_base = square_top - side - 12 - S(_u("stepName"))
            pdf.drawString(item_left, name_base, name)
            pdf.setFillColorRGB(*muted)
            text_size = S(_u("stepText"))
            text_lead = text_size * deck_type.LEADING["stepText"]
            pdf.setFont(font, text_size)
            for offset, line in enumerate(_wrap(text, font, text_size, span - 4)[:4]):
                pdf.drawString(item_left, name_base - 8 - text_size - offset * text_lead, line)
        return

    if layout == "cards":
        gap = 18.0
        span = (width - gap * (len(pairs) - 1)) / max(len(pairs), 1)
        height = min((60 if compact else 100) * _K, room - 10)
        card_top = top - 10
        name_size = S(14.0 if compact else _u("cardName"))
        text_size = S(12.0 if compact else _u("cardText"))
        text_lead = text_size * deck_type.LEADING["cardText"]
        for index, (name, text) in enumerate(pairs):
            item_left = left + index * (span + gap)
            _pdf_box(
                pdf,
                look,
                left=item_left,
                bottom=card_top - height,
                width=span,
                height=height,
                fill=tint,
                line=hair,
                bg=bg,
            )
            pdf.setFillColorRGB(*accent)
            pdf.rect(item_left, card_top - 5, span, 5, stroke=0, fill=1)
            pdf.setFont(bold, name_size)
            name_lines = _wrap(name, bold, name_size, span - 28)[:2]
            name_base = card_top - 20 - name_size
            for offset, line in enumerate(name_lines):
                pdf.drawString(item_left + 14, name_base - offset * name_size * 1.3, line)
            name_base -= (len(name_lines) - 1) * name_size * 1.3
            pdf.setFillColorRGB(*ink)
            pdf.setFont(font, text_size)
            first = name_base - 12 - text_size
            lines = _wrap(text, font, text_size, span - 28)[
                : max(1, int((first - (card_top - height) - 10) / text_lead) + 1)
            ]
            for offset, line in enumerate(lines):
                pdf.drawString(item_left + 14, first - offset * text_lead, line)
        return

    # timeline
    axis = 128.0
    step = min(56.0, room / max(len(pairs), 1))
    line_size = S(max(_u("lineMin"), min(_u("lineMax"), step / 2.3)))
    line_lead = line_size * deck_type.LEADING["line"]
    pdf.setFillColorRGB(*tint)
    pdf.rect(left + axis, top - step * len(pairs), 1.5, step * len(pairs), stroke=0, fill=1)
    for index, (when, what) in enumerate(pairs):
        line_top = top - 4 - line_size - index * step
        pdf.setFillColorRGB(*accent)
        pdf.setFont(bold, line_size)
        pdf.drawRightString(left + axis - 12, line_top, when)
        pdf.rect(left + axis - 3.25, line_top - 1, 8, 8, stroke=0, fill=1)
        pdf.setFillColorRGB(*ink)
        pdf.setFont(font, line_size)
        for offset, line in enumerate(_wrap(what, font, line_size, width - axis - 16)[:2]):
            pdf.drawString(left + axis + 16, line_top - offset * line_lead, line)


def _pdf_chart(
    pdf,
    chart: dict,
    *,
    accent: tuple[float, float, float],
    muted: tuple[float, float, float],
    top: float,
    width: float,
    font: str,
    left: float = 72.0,
) -> None:
    """The `.pdf` twin of `_pptx_chart`, drawn by hand over a zero floor."""
    bottom = _H - _BODY_BOTTOM + 6
    height = top - bottom - 30
    if height < 60 or width < 100:
        return

    values = [value for _, series in chart["series"] for value in series]
    if max(values + [0]) <= 0:
        return
    ceiling = _nice_ceiling(max(values))
    categories = chart["categories"]
    step = width / len(categories)

    pdf.setFont(font, 11)
    for tick in range(5):
        y = bottom + height * tick / 4
        pdf.setStrokeColorRGB(0.9, 0.9, 0.9)
        pdf.setLineWidth(0.5)
        pdf.line(left, y, left + width, y)
        pdf.setFillColorRGB(*muted)
        pdf.drawRightString(left - 6, y - 4, _tick_label(ceiling * tick / 4))
    if unit := chart.get("unit"):
        pdf.setFillColorRGB(*muted)
        pdf.drawString(left - 30, bottom + height + 12, unit)

    for series_index, (_name, series) in enumerate(chart["series"]):
        colour = accent if series_index == 0 else tuple(c + (1 - c) * 0.55 for c in accent)
        if chart["kind"] == "line":
            pdf.setStrokeColorRGB(*colour)
            pdf.setLineWidth(2.5)
            path = pdf.beginPath()
            points = [
                (left + step * (position + 0.5), bottom + height * (value / ceiling))
                for position, value in enumerate(series)
            ]
            for position, (x, y) in enumerate(points):
                path.moveTo(x, y) if position == 0 else path.lineTo(x, y)
            pdf.drawPath(path)
            pdf.setFillColorRGB(*colour)
            for x, y in points:
                pdf.circle(x, y, 3, stroke=0, fill=1)
        else:
            # Series share the category slot side by side.
            count = len(chart["series"])
            span = step * 0.6 / count
            pdf.setFillColorRGB(*colour)
            for position, value in enumerate(series):
                x = left + step * (position + 0.5) - (step * 0.6) / 2 + span * series_index
                pdf.rect(x, bottom, span, height * (value / ceiling), stroke=0, fill=1)

    pdf.setFillColorRGB(*muted)
    pdf.setFont(font, 12)
    for position, label in enumerate(categories):
        pdf.drawCentredString(left + step * (position + 0.5), bottom - 18, label)

    if len(chart["series"]) > 1:
        x = left
        pdf.setFont(font, 12)
        for series_index, (name, _values) in enumerate(chart["series"]):
            colour = accent if series_index == 0 else tuple(c + (1 - c) * 0.55 for c in accent)
            pdf.setFillColorRGB(*colour)
            pdf.circle(x + 4, bottom - 36, 3.5, stroke=0, fill=1)
            pdf.setFillColorRGB(*muted)
            pdf.drawString(x + 13, bottom - 40, name)
            x += 24 + pdfmetrics.stringWidth(name, font, 12)


def _nice_ceiling(highest: float) -> float:
    """Top of scale rounded up to 1, 2, 2.5, 5 or 10 × a power of ten, as PowerPoint does."""
    import math

    if highest <= 0:
        return 1.0
    power = 10 ** math.floor(math.log10(highest))
    for multiple in (1, 2, 2.5, 5, 10):
        if highest <= multiple * power:
            return multiple * power
    return highest


def _tick_label(value: float) -> str:
    """A gridline's number: integer at 10 and above, one decimal below."""
    return f"{value:,.0f}" if abs(value) >= 10 else f"{value:,.1f}".rstrip("0").rstrip(".")


def _typescale(slide: dict) -> float:
    """The slide's `textScale` clamped to 0.5–2.0; it arrives on a PATCHable artifact."""
    try:
        value = float(slide.get("textScale") or 1.0)
    except (TypeError, ValueError):
        return 1.0
    return min(2.0, max(0.5, value))


def _wrap(text: str, font: str, size: float, width: float) -> list[str]:
    """Greedy wrap by measured width, breaking between characters (Korean has no spaces)."""
    lines: list[str] = []
    current = ""
    for char in text:
        if char == "\n":
            lines.append(current)
            current = ""
            continue
        candidate = current + char
        if pdfmetrics.stringWidth(candidate, font, size) > width and current:
            # Prefer the last space, so Latin text still breaks on words.
            cut = current.rfind(" ")
            if cut > len(current) * 0.6:
                lines.append(current[:cut])
                current = (current[cut + 1 :] + char).lstrip(" ")
            else:
                lines.append(current.rstrip(" "))
                current = char.lstrip(" ")
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines


def to_pdf(title: str, slides: list[dict], *, tokens: dict[str, str] | None = None) -> bytes:
    """The deck as a PDF, one slide per page, without notes."""
    style = design.normalise_tokens(tokens) if tokens else None
    font = fonts.korean(style["font"] if style else "gothic")
    bold = fonts.korean(style["font"] if style else "gothic", bold=True)
    ink = _hex_floats(style["ink"]) if style else _PDF_INK
    muted = _hex_floats(style["muted"]) if style else _PDF_MUTED
    mark = _logo_of(style)
    footer = (style or {}).get("footer") or ""
    visual_style = (style or {}).get("visualStyle") or "editorial"
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=(_W, _H))
    pdf.setTitle(title)

    look = _look_of(visual_style)
    if look.own_neutrals:
        ink, muted = _hex_floats(look.ink), _hex_floats(look.muted)
    ground = _hex_floats(look.bg)
    for index, data in enumerate(_written(slides)):
        accent = _hex_floats(data.get("accent") or (style or {}).get("accent"))
        if visual_style == "mono":
            accent = ink
        layout = data.get("layout") or "bullets"
        heading = str(data.get("title") or "")
        body = str(data.get("body") or "")
        bullets = [str(b) for b in (data.get("bullets") or []) if str(b).strip()]
        rows = [[str(cell) for cell in row] for row in (data.get("rows") or []) if row]
        metrics = [
            [str(pair[0]), str(pair[1])]
            for pair in (data.get("metrics") or [])
            if isinstance(pair, list) and len(pair) >= 2
        ]
        chart = _chart_of(data)
        picture = _picture_of(data)
        picture_span = _picture_span(data)
        pairs = _pairs_of(data, layout)
        picture_left = bool(
            picture
            and (bullets or rows or metrics or chart or body)
            and str((data.get("image") or {}).get("position") or "") == "left"
        )
        # Same split as the .pptx.
        text_width = (
            _W
            - 144
            - (
                picture_span + 24
                if picture and (bullets or rows or metrics or chart or body)
                else 0
            )
        )
        text_left = 72 + (picture_span + 24 if picture_left else 0)

        cover = layout in _COVERS
        on_accent = look.cover in ("gradient", "glow")
        cover_left = 72 + 336 if look.cover == "split" else 72
        tint = _mix_floats(accent, look.tint, onto=ground) if look.tint else (0.949, 0.949, 0.949)
        hair = _hex_floats(look.hair)
        if cover:
            if look.cover == "gradient":
                pdf.setFillColorRGB(*accent)
                pdf.rect(0, 0, _W, _H, stroke=0, fill=1)
                pdf.saveState()
                pdf.linearGradient(
                    0,
                    _H,
                    _W,
                    0,
                    [
                        Color(*accent),
                        Color(
                            *_mix_floats(
                                accent,
                                48 if visual_style == "poster" else 62,
                                onto=_PDF_INK,
                            )
                        ),
                    ],
                    extend=True,
                )
                pdf.restoreState()
            elif look.cover == "wash":
                pdf.setFillColorRGB(*_mix_floats(accent, 10))
                pdf.rect(0, 0, _W, _H, stroke=0, fill=1)
            else:
                pdf.setFillColorRGB(*ground)
                pdf.rect(0, 0, _W, _H, stroke=0, fill=1)
                if look.cover == "glow":
                    for ring, pct in ((520, 22), (400, 38), (280, 58)):
                        pdf.setFillColorRGB(*_mix_floats(accent, pct, onto=ground))
                        pdf.circle(_W - ring * 0.05, ring * 0.05, ring / 2, stroke=0, fill=1)
                elif look.cover == "paper":
                    pdf.setFillColorRGB(*accent)
                    pdf.circle(_W - 72, _H - 272, 228, stroke=0, fill=1)
                elif look.cover == "split":
                    pdf.setFillColorRGB(*accent)
                    pdf.rect(0, 0, 384, _H, stroke=0, fill=1)
                elif look.cover == "brackets":
                    pdf.setFillColorRGB(*ink)
                    pdf.rect(62, _H - 139, 6, 72, stroke=0, fill=1)
                    pdf.rect(62, _H - 73, 53, 6, stroke=0, fill=1)
                    pdf.rect(_W - 68, 67, 6, 72, stroke=0, fill=1)
                    pdf.rect(_W - 115, 67, 53, 6, stroke=0, fill=1)
        else:
            pdf.setFillColorRGB(*ground)
            pdf.rect(0, 0, _W, _H, stroke=0, fill=1)
            if look.ornament == "left-bar":
                pdf.setFillColorRGB(*accent)
                pdf.rect(0, 0, 14, _H, stroke=0, fill=1)
            elif look.ornament == "top-band":
                pdf.setFillColorRGB(*accent)
                pdf.rect(0, _H - 14, _W, 14, stroke=0, fill=1)
            elif look.ornament == "corner-circle":
                pdf.setFillColorRGB(*_mix_floats(accent, 12, onto=ground))
                pdf.circle(_W - 12, _H, 84, stroke=0, fill=1)
            elif look.ornament == "bottom-rule":
                pdf.setFillColorRGB(*accent)
                pdf.rect(0, 0, _W * 0.6, 10, stroke=0, fill=1)
                pdf.setFillColorRGB(*_mix_floats(accent, 40, onto=ground))
                pdf.rect(_W * 0.6, 0, _W * 0.4, 10, stroke=0, fill=1)
            elif look.ornament == "gutter":
                pdf.setFillColorRGB(*accent)
                pdf.rect(0, 0, 7, _H, stroke=0, fill=1)
                pdf.setFont(bold, _u("gutterNumber"))
                pdf.drawString(22, 80, f"{index + 1:02d}")
            elif look.ornament == "frame":
                pdf.setStrokeColorRGB(*ink)
                pdf.setLineWidth(1.0)
                pdf.rect(24, 24, _W - 48, _H - 48, stroke=1, fill=0)
            elif look.ornament == "bottom-band":
                pdf.setFillColorRGB(*tint)
                pdf.rect(0, 0, _W, 53, stroke=0, fill=1)

        # `textScale` applies to sizes and line advances alike: this canvas
        # does not reflow.
        ts = _typescale(data)

        def S(n: float, _ts: float = ts) -> float:
            # A size at the slide's scale, never under the floor.
            return max(float(deck_type.FLOOR_PT), n * _ts)

        cover_size = _u(
            "coverPoster"
            if visual_style == "poster"
            else "coverMono"
            if visual_style == "mono"
            else "cover"
        )
        #: Body box for this slide; a long title lowers its top.
        body_top, room = _BODY_TOP, _BODY_BOTTOM - _BODY_TOP

        if layout == "closing":
            cover_ink = (1.0, 1.0, 1.0) if on_accent else ink
            cover_muted = _mix_floats((1.0, 1.0, 1.0), 80, onto=accent) if on_accent else muted
            if look.cover == "split":
                pdf.setFillColorRGB(*_mix_floats((1.0, 1.0, 1.0), 35, onto=accent))
                pdf.setFont(bold, S(_u("splitNumber")))
                pdf.drawString(60, 60, "END")
            if look.cover != "brackets":
                pdf.setFillColorRGB(*((1, 1, 1) if on_accent else accent))
                pdf.rect(cover_left, _H - 127, 106, 7, stroke=0, fill=1)
            pdf.setFillColorRGB(*cover_ink)
            closing_size = S(_u("closing"))
            pdf.setFont(bold, closing_size)
            y = _H - 140 - closing_size
            for line in _wrap(heading or "마무리", bold, closing_size, _W - 72 - cover_left)[:2]:
                pdf.drawString(cover_left, y, line)
                y -= closing_size * 1.2
            y -= 6
            bullet_size = S(_u("closingBullets"))
            bullet_lead = bullet_size * deck_type.LEADING["body"]
            for text in bullets[:3]:
                pdf.setFillColorRGB(*cover_muted)
                pdf.setFont(font, bullet_size)
                pdf.drawString(cover_left, y, "—")
                pdf.setFillColorRGB(*cover_ink)
                for offset, line in enumerate(
                    _wrap(text, font, bullet_size, _W - 128 - cover_left)[:2]
                ):
                    pdf.drawString(cover_left + 28, y - offset * bullet_lead, line)
                    y -= bullet_lead
                y -= 6
            if body:
                pdf.setFillColorRGB(*cover_ink)
                pdf.setFont(bold, S(_u("closingBody")))
                pdf.drawString(
                    cover_left, 92, _wrap(body, bold, S(_u("closingBody")), _W - 72 - cover_left)[0]
                )
        elif cover:
            cover_ink = (1.0, 1.0, 1.0) if on_accent else ink
            number = str(data.get("number") or "")
            if look.cover == "split":
                pdf.setFillColorRGB(*_mix_floats((1.0, 1.0, 1.0), 35, onto=accent))
                pdf.setFont(bold, S(_u("splitNumber")))
                pdf.drawString(60, 60, number.replace(".", "") or "01")
            if look.cover != "split" and layout == "section" and number:
                pdf.setFillColorRGB(
                    *(_mix_floats((1.0, 1.0, 1.0), 70, onto=accent) if on_accent else accent)
                )
                pdf.setFont(bold, S(_u("sectionNumber")))
                pdf.drawString(cover_left, _H / 2 + 100, number)
            elif look.cover != "brackets":
                pdf.setFillColorRGB(*((1, 1, 1) if on_accent else accent))
                pdf.rect(cover_left + 10, _H / 2 + 74, 106, 7, stroke=0, fill=1)
            pdf.setFillColorRGB(*cover_ink)
            pdf.setFont(bold, S(cover_size))
            y = _H / 2 + 20
            title_width = (_W - 72 - cover_left) - (300 if look.cover == "paper" else 0)
            for line in _wrap(heading or title, bold, S(cover_size), title_width):
                pdf.drawString(cover_left, y, line)
                y -= S(cover_size) * 1.2
            if body:
                pdf.setFillColorRGB(
                    *(_mix_floats((1.0, 1.0, 1.0), 80, onto=accent) if on_accent else muted)
                )
                pdf.setFont(font, S(_u("coverBody")))
                for line in _wrap(body, font, S(_u("coverBody")), _W - 72 - cover_left)[:2]:
                    pdf.drawString(cover_left, y - 6, line)
                    y -= S(_u("coverBody")) * 1.5
        elif layout == "statement":
            pdf.setFillColorRGB(*accent)
            pdf.rect((_W - 62) / 2, _H - 181, 62, 5, stroke=0, fill=1)
            pdf.setFont(bold, S(_u("statement")))
            y = _H / 2 + 10
            for line in _wrap(heading, bold, S(_u("statement")), _W - 180)[:2]:
                pdf.drawCentredString(_W / 2, y, line)
                y -= S(_u("statement")) * 1.25
            if body:
                pdf.setFillColorRGB(*muted)
                pdf.setFont(font, S(_u("statementBody")))
                for line in _wrap(body, font, S(_u("statementBody")), _W - 240)[:2]:
                    pdf.drawCentredString(_W / 2, y - 4, line)
                    y -= S(_u("statementBody")) * 1.5
        elif layout == "quote":
            pdf.setFillColorRGB(*accent)
            pdf.setFont(bold, S(_u("quote")))
            y = _H / 2 + 40
            for line in _wrap(f"“{body or heading}”", bold, S(_u("quote")), _W - 200):
                pdf.drawString(90, y, line)
                y -= S(_u("quote")) * 1.4
            if body and heading:
                pdf.setFillColorRGB(*muted)
                pdf.setFont(font, S(_u("quoteBy")))
                pdf.drawString(90, y - 8, heading)
        else:
            title_size = deck_type.title_pt(heading, text_width / _K)
            title_line = title_size * deck_type.LEADING["title"]
            title_lines = _wrap(heading, bold, title_size, text_width) or [""]
            pdf.setFillColorRGB(*ink)
            pdf.setFont(bold, title_size)
            y = _H - _TITLE_TOP - title_size * 0.9
            for line in title_lines:
                pdf.drawString(text_left, y, line)
                y -= title_line
            # The tab under the title; the body box starts below it and loses one line for
            # each extra title line.
            pdf.setFillColorRGB(*accent)
            tab_top = _TITLE_TOP + title_line * len(title_lines) + 6 * _K
            pdf.rect(text_left, _H - tab_top - 5, 62, 5, stroke=0, fill=1)
            body_top = _BODY_TOP + title_line * (len(title_lines) - 1)
            room = _BODY_BOTTOM - body_top

            # The figure band, as in the pptx renderer.
            compact = False
            if picture and (data.get("image") or {}).get("diagram") and _has_words(data):
                image_bytes, _caption = picture
                text_width, text_left = _W - 144, 72.0
                band_width, band_height = _fit(image_bytes, box=(text_width, room * 0.56))
                try:
                    pdf.drawImage(
                        ImageReader(io.BytesIO(image_bytes)),
                        text_left + (text_width - band_width) / 2,
                        _H - body_top - band_height,
                        width=band_width,
                        height=band_height,
                        mask="auto",
                    )
                except Exception as exc:  # noqa: BLE001 — a bad picture is not a failed export
                    log.warning("could not place a figure band into the deck pdf: %s", exc)
                body_top += band_height + 8
                room -= band_height + 8
                picture = None
                compact = True
            y = _H - body_top

            if pairs:
                _pdf_pairs(
                    pdf,
                    pairs,
                    layout=layout,
                    accent=accent,
                    tint=tint,
                    muted=muted,
                    ink=ink,
                    top=y,
                    width=text_width,
                    font=font,
                    scale=ts,
                    left=text_left,
                    look=look,
                    hair=hair,
                    bg=ground,
                    room=room,
                    bold=bold,
                    compact=compact,
                )
            elif chart:
                _pdf_chart(
                    pdf,
                    chart,
                    accent=accent,
                    muted=muted,
                    top=y,
                    width=text_width,
                    font=font,
                    left=text_left,
                )
            elif metrics and layout == "big-number":
                figure, label = metrics[0]
                figure_size = S(_u("bigNumber"))
                pdf.setFillColorRGB(*accent)
                pdf.setFont(bold, figure_size)
                pdf.drawString(text_left, y - figure_size, figure)
                figure_width = pdf.stringWidth(figure, font, figure_size)
                pdf.setFillColorRGB(*muted)
                pdf.setFont(font, S(_u("bigNumberLabel")))
                pdf.drawString(text_left + figure_width + 14, y - figure_size, label)
                if body:
                    body_size = S(_u("bigNumberBody"))
                    pdf.setFillColorRGB(*ink)
                    pdf.setFont(font, body_size)
                    for offset, line in enumerate(_wrap(body, font, body_size, text_width)[:2]):
                        pdf.drawString(
                            text_left,
                            y - figure_size - 24 - body_size - offset * body_size * 1.5,
                            line,
                        )
            elif metrics:
                span = (text_width - 24 * (len(metrics) - 1)) / len(metrics)
                figure_size = S(_u("metric"))
                label_size = S(_u("metricLabel"))
                box_height = figure_size * 1.1 + 6 + label_size * 1.5 + 30
                for position, (figure, label) in enumerate(metrics):
                    left = text_left + position * (span + 24)
                    _pdf_box(
                        pdf,
                        look,
                        left=left,
                        bottom=y - 8 - box_height,
                        width=span,
                        height=box_height,
                        fill=tint,
                        line=hair,
                        bg=ground,
                    )
                    pdf.setFillColorRGB(*accent)
                    pdf.rect(left, y - 8 - 5, span, 5, stroke=0, fill=1)
                    pdf.setFillColorRGB(*accent)
                    pdf.setFont(bold, figure_size)
                    pdf.drawString(left + 14, y - 8 - 14 - figure_size, figure)
                    pdf.setFillColorRGB(*muted)
                    pdf.setFont(font, label_size)
                    pdf.drawString(
                        left + 14, y - 8 - 14 - figure_size * 1.1 - 6 - label_size, label
                    )
            elif rows:
                # Columns by their widest cell and rows by the shared row height, as in the
                # panel; a cell that wraps makes its row taller, up to two lines.
                widths = [text_width * share for share in deck_type.column_shares(rows)]
                cell_size = S(deck_type.table_size(len(rows)) * _K)
                cell_pad = S(deck_type.table_pad(len(rows)) * _K)
                line_height = cell_size * deck_type.LEADING["table"]
                pad_x = 9 * _K
                row_top = y
                for row_index, row in enumerate(rows):
                    wrapped = [
                        _wrap(
                            cell,
                            bold if row_index == 0 or c == 0 else font,
                            cell_size,
                            widths[c] - 2 * pad_x,
                        )[:2]
                        for c, cell in enumerate(row)
                        if c < len(widths)
                    ]
                    step = max(1, *(len(w) for w in wrapped)) * line_height + 2 * cell_pad
                    if row_top - step < _H - _BODY_BOTTOM - 2:
                        break
                    if row_index == 0:
                        pdf.setFillColorRGB(*accent)
                        pdf.rect(text_left, row_top - step, text_width, step, stroke=0, fill=1)
                    elif row_index % 2 == 0:
                        pdf.setFillColorRGB(*tint)
                        pdf.rect(text_left, row_top - step, text_width, step, stroke=0, fill=1)
                    else:
                        pdf.setStrokeColorRGB(*hair)
                        pdf.setLineWidth(0.75)
                        pdf.line(text_left, row_top - step, text_left + text_width, row_top - step)
                    pdf.setFillColorRGB(*((1.0, 1.0, 1.0) if row_index == 0 else ink))
                    x = text_left
                    for c, lines in enumerate(wrapped):
                        pdf.setFont(bold if row_index == 0 or c == 0 else font, cell_size)
                        baseline = row_top - cell_pad - cell_size * 0.85
                        for offset, line in enumerate(lines):
                            pdf.drawString(x + pad_x, baseline - offset * line_height, line)
                        x += widths[c]
                    row_top -= step
            elif bullets and layout == "agenda":
                columns = _split_columns(bullets) if len(bullets) > 4 else [bullets]
                span = (text_width - 24 * (len(columns) - 1)) / len(columns)
                per = max(len(column) for column in columns)
                step = min(66.0, room / max(per, 1))
                agenda_size = S(_u("agenda"))
                number = 0
                top = y
                for column_index, column in enumerate(columns):
                    left = text_left + column_index * (span + 24)
                    line_y = top - 10
                    for text in column:
                        number += 1
                        pdf.setFillColorRGB(*accent)
                        pdf.setFont(bold, S(_u("agendaNumber")))
                        pdf.drawString(left, line_y - agenda_size, f"{number:02d}")
                        pdf.setFillColorRGB(*ink)
                        pdf.setFont(font, agenda_size)
                        lines = _wrap(text, font, agenda_size, span - 64)[:2]
                        for offset, line in enumerate(lines):
                            pdf.drawString(
                                left + 64, line_y - agenda_size - offset * agenda_size * 1.3, line
                            )
                        # The rule under the entry; a wrapped entry pushes the next one down.
                        rule_y = line_y - agenda_size - (len(lines) - 1) * agenda_size * 1.3 - 12
                        pdf.setFillColorRGB(*hair)
                        pdf.rect(left, rule_y, span, 0.75, stroke=0, fill=1)
                        line_y = min(line_y - step, rule_y - 10)
            elif bullets:
                columns = _columns_of(data, bullets, layout)
                size = S(_u("bodyNarrow") if len(columns) > 1 else _u("body"))
                step = size * look.leading
                gap = size * deck_type.BULLET_GAP
                span = (text_width - 24 * (len(columns) - 1)) / len(columns)
                wrapped_columns = [
                    [_wrap(text, font, size, span - 26) for text in column] for column in columns
                ]
                top = y - size * 0.85
                for column_index, column in enumerate(wrapped_columns):
                    left = text_left + column_index * (span + 24)
                    y = top
                    for wrapped in column:
                        pdf.setFillColorRGB(*accent)
                        pdf.setFont(font, size)
                        pdf.drawString(left + 4, y, "•")
                        pdf.setFillColorRGB(*ink)
                        for offset, line in enumerate(wrapped):
                            pdf.drawString(left + 26, y - offset * step, line)
                        y -= step * len(wrapped) + gap
            elif body:
                size = S(_u("paragraph"))
                step = size * look.leading
                wrapped = _wrap(body, font, size, text_width)
                y = y - size * 0.85
                pdf.setFillColorRGB(*muted)
                pdf.setFont(font, size)
                for line in wrapped:
                    pdf.drawString(text_left, y, line)
                    y -= step

        if picture:
            image_bytes, image_caption = picture
            alone = not (bullets or rows or metrics or chart or body)
            box = (_W - 260, room) if alone else (picture_span, room)
            fill = str((data.get("image") or {}).get("fit") or "") == "cover"
            if fill:
                width, height, *_ = _fill(image_bytes, box=box)
                box_width, box_height = box
                box_left = (
                    (72 if picture_left else 72 + text_width + 24)
                    if not alone
                    else (_W - box_width) / 2
                )
                box_bottom = _H - body_top - room
                left = box_left - (width - box_width) / 2
                bottom = box_bottom - (height - box_height) / 2
            else:
                width, height = _fit(image_bytes, box=box)
                box_width, box_height = width, height
                box_left = (
                    (72 if picture_left else 72 + text_width + 24)
                    if not alone
                    else (_W - width) / 2
                )
                box_bottom = _H - body_top - room + max(0.0, (room - height) / 2)
                left, bottom = box_left, box_bottom
            try:
                if fill:
                    pdf.saveState()
                    clip = pdf.beginPath()
                    clip.rect(box_left, box_bottom, box_width, box_height)
                    pdf.clipPath(clip, stroke=0, fill=0)
                pdf.drawImage(
                    ImageReader(io.BytesIO(image_bytes)),
                    left,
                    bottom,
                    width=width,
                    height=height,
                    mask="auto",
                )
                if fill:
                    pdf.restoreState()
            except Exception as exc:  # noqa: BLE001 — a bad picture is not a failed export
                if fill:
                    try:
                        pdf.restoreState()
                    except Exception:  # noqa: BLE001 — recovery only
                        pass
                log.warning("could not draw a picture into the deck pdf: %s", exc)
            else:
                if image_caption:
                    pdf.setFillColorRGB(*muted)
                    pdf.setFont(font, S(_u("caption")))
                    pdf.drawString(
                        box_left,
                        box_bottom - 6 - S(_u("caption")),
                        _wrap(image_caption, font, S(_u("caption")), box_width)[0],
                    )

        # Foot: logo, deck title and footer on the left, slide number on the right.
        if not cover:
            pdf.setStrokeColorRGB(*hair)
            pdf.setLineWidth(0.75)
            pdf.line(72, _FOOT_RULE, _W - 72, _FOOT_RULE)
            left = 72.0
            if mark:
                blob, mark_width, mark_height = mark
                try:
                    pdf.drawImage(
                        ImageReader(io.BytesIO(blob)),
                        left,
                        30,
                        width=mark_width,
                        height=mark_height,
                        mask="auto",
                    )
                except Exception as exc:  # noqa: BLE001 — a mark, not the deck
                    log.warning("could not draw the logo on a pdf slide: %s", exc)
                else:
                    left += mark_width + 10
            foot_size = _u("footer")
            pdf.setFillColorRGB(0.55, 0.55, 0.55)
            pdf.setFont(font, foot_size)
            foot = _wrap(title, font, foot_size, _W - 400 - (left - 72))[0] if title else ""
            pdf.drawString(left, 36, foot)
            if footer:
                pdf.drawRightString(_W - 120, 36, _wrap(footer, font, foot_size, 240)[0])
            pdf.setFillColorRGB(*accent)
            pdf.rect(_W - 108, 24, 36, 36, stroke=0, fill=1)
            pdf.setFillColorRGB(1, 1, 1)
            pdf.setFont(bold, _u("pageNumber"))
            pdf.drawCentredString(_W - 90, 24 + 18 - _u("pageNumber") * 0.35, str(index + 1))
        pdf.showPage()

    pdf.save()
    return buffer.getvalue()


__all__ = ["to_pdf", "to_pptx"]
