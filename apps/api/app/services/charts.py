"""DrawingML chart styling shared by the deck (.pptx) and the report (.docx).

`apply` styles an existing python-pptx chart; `part` builds one through a
throwaway presentation and returns the chart part and workbook blobs a `.docx`
embeds.
"""

from __future__ import annotations

import logging

log = logging.getLogger(__name__)

#: (latin, East Asian) pairs. `font.name` sets `a:latin` only; Hangul is drawn
#: from `a:ea`.
FACES = {
    "gothic": ("Segoe UI", "맑은 고딕"),
    "serif": ("Georgia", "바탕"),
}

_GRID = "E5E5E5"
_AXIS = "C8C8C8"


def apply(chart, *, kind: str, unit: str, accent, muted, faces: tuple[str, str]) -> None:
    """Styles a chart: zero-floored value axis, light gridlines, accent series, document faces."""
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LEGEND_POSITION, XL_MARKER_STYLE
    from pptx.util import Pt

    _plain_frame(chart)
    chart.has_title = False
    chart.has_legend = len(chart.plots[0].series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        _text(chart.legend.font, muted, 12, faces)

    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor.from_string(_GRID)
    value_axis.format.line.color.rgb = RGBColor.from_string(_GRID)
    _text(value_axis.tick_labels.font, muted, 12, faces)
    if unit:
        value_axis.has_title = True
        frame = value_axis.axis_title.text_frame
        frame.text = unit
        _text(frame.paragraphs[0].font, muted, 12, faces)
        # Upright: the default rotated title lays a Hangul unit on its side.
        frame._txBody.bodyPr.set("rot", "0")
        frame._txBody.bodyPr.set("vert", "horz")

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.format.line.color.rgb = RGBColor.from_string(_AXIS)
    _text(category_axis.tick_labels.font, muted, 12, faces)

    # First series in the accent, later ones in a lighter mix of it.
    for position, series in enumerate(chart.plots[0].series):
        colour = accent if position == 0 else mix(accent, RGBColor(0xFF, 0xFF, 0xFF), 0.55)
        if kind == "line":
            series.format.line.color.rgb = colour
            series.format.line.width = Pt(2.5)
            # Markers otherwise keep the theme palette.
            series.marker.style = XL_MARKER_STYLE.CIRCLE
            series.marker.size = 6
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = colour
            series.marker.format.line.color.rgb = colour
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colour
            series.format.line.fill.background()


def _plain_frame(chart) -> None:
    """Removes the frame and rounded corners Word draws around an unstyled chart.

    `c:spPr` is schema-ordered and must follow `c:chart`, or Word offers to
    repair the file.
    """
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn

    space = chart._chartSpace
    for tag in ("c:roundedCorners", "c:spPr"):
        for existing in space.findall(qn(tag)):
            space.remove(existing)

    corners = parse_xml(f'<c:roundedCorners {nsdecls("c")} val="0"/>')
    space.insert(0, corners)

    plot = space.find(qn("c:chart"))
    shape = parse_xml(f"<c:spPr {nsdecls('c', 'a')}><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>")
    space.insert(list(space).index(plot) + 1, shape)


def part(
    kind: str,
    categories: list[str],
    series: list[tuple[str, list[float]]],
    *,
    unit: str,
    accent,
    muted,
    faces: tuple[str, str],
) -> tuple[bytes, bytes] | None:
    """`(chart part blob, workbook blob)` built via a throwaway presentation, or None on failure."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu

    try:
        payload = CategoryChartData()
        payload.categories = categories
        for name, values in series:
            payload.add_series(name or " ", values)

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS if kind == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(0),
            Emu(0),
            Emu(5486400),
            Emu(3200400),
            payload,
        )
        apply(frame.chart, kind=kind, unit=unit, accent=accent, muted=muted, faces=faces)
        chart_part = frame.chart.part
        return chart_part.blob, chart_part.chart_workbook.xlsx_part.blob
    except Exception as exc:  # noqa: BLE001 — a chart is not worth a failed export
        log.warning("could not build a chart part: %s", exc)
        return None


def mix(colour, toward, amount: float):
    """`colour` moved `amount` (0..1) toward `toward`."""
    from pptx.dml.color import RGBColor

    return RGBColor(*(round(a + (b - a) * amount) for a, b in zip(colour, toward, strict=True)))


def _text(font, colour, size: int, faces: tuple[str, str]) -> None:
    """Sets size, colour and both latin and East Asian faces on a chart font."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    font.size = Pt(size)
    font.color.rgb = colour
    font.name = faces[0]
    properties = font._rPr
    for tag in ("a:ea", "a:cs"):
        properties.append(properties.makeelement(qn(tag), {"typeface": faces[1]}))


__all__ = ["FACES", "apply", "mix", "part"]
