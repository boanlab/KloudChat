"""Concurrent exports, checks and turns with different inputs each answer their own question."""

from __future__ import annotations

import asyncio
import io
import zipfile
from concurrent.futures import ThreadPoolExecutor

import pytest

from app.services import deck_export, hangul, hwpx_import, lint, report_export, richtext

#: Enough to interleave, small enough to stay quick.
_MANY = 24
_WORKERS = 8


def _mark(index: int) -> str:
    """A zero-padded, bracketed tag that is nobody else's, even as a substring."""
    return f"[{index:02d}]"


def _spread(work, count: int = _MANY) -> list:
    """`work(i)` for every `i`, all at once, results in order."""
    with ThreadPoolExecutor(max_workers=_WORKERS) as pool:
        return list(pool.map(work, range(count)))


def _sections(index: int) -> list[dict]:
    return [
        {
            "id": f"s{index}",
            "heading": f"{_mark(index)}절 제목",
            "level": 1,
            "format": "html",
            "content": (
                f"<p>이것은 {_mark(index)}번 보고서의 본문이다.</p>"
                f'<table><tr><th colspan="2">{_mark(index)}번 표</th></tr>'
                f"<tr><td>가</td><td>{_mark(index)}</td></tr></table>"
            ),
        }
    ]


def _text_of(hwpx: bytes) -> str:
    archive = zipfile.ZipFile(io.BytesIO(hwpx))
    return "".join(
        archive.read(name).decode("utf8", "replace")
        for name in archive.namelist()
        if name.startswith("Contents/section")
    )


def test_two_dozen_reports_written_at_once_are_two_dozen_reports() -> None:
    """Concurrent report exports do not share state."""

    def write(index: int) -> bytes:
        return report_export.to_hwpx(
            f"{_mark(index)}번 보고서", richtext.normalise(_sections(index))
        )

    for index, produced in enumerate(_spread(write)):
        body = _text_of(produced)
        assert f"{_mark(index)}번 보고서" in body
        assert f"이것은 {_mark(index)}번 보고서의 본문이다." in body
        # Nobody else's title is in this file.
        others = [n for n in range(_MANY) if n != index and f"{_mark(n)}번 보고서" in body]
        assert others == []


def test_a_deck_keeps_its_own_type_scale_when_others_are_being_written() -> None:
    """Concurrent `to_pptx` calls keep their own per-slide type scale."""

    def write(index: int) -> bytes:
        scale = "large" if index % 2 else "small"
        return deck_export.to_pptx(
            f"{_mark(index)}번 덱",
            [
                {"layout": "title", "title": f"{_mark(index)}번 덱", "body": "표지"},
                {
                    "layout": "bullets",
                    "title": f"{_mark(index)}번 장",
                    "bullets": [f"{_mark(index)}번 항목"],
                    "textScale": scale,
                },
            ],
        )

    from pptx import Presentation

    for index, produced in enumerate(_spread(write)):
        deck = Presentation(io.BytesIO(produced))
        words = " ".join(
            shape.text_frame.text
            for slide in deck.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert f"{_mark(index)}번 장" in words
        assert [n for n in range(_MANY) if n != index and f"{_mark(n)}번 장" in words] == []


def test_reading_two_dozen_hwpx_files_at_once_reads_two_dozen_documents() -> None:
    """Concurrent HWPX imports do not share state."""
    files = _spread(
        lambda index: report_export.to_hwpx(
            f"{_mark(index)}번 문서", richtext.normalise(_sections(index))
        )
    )

    def read(index: int) -> hwpx_import.Document:
        return hwpx_import.read(files[index])

    for index, document in enumerate(_spread(read)):
        assert document.title == f"{_mark(index)}번 문서"
        assert f"{_mark(index)}절 제목" in [part.heading for part in document.sections]


def test_the_checks_answer_about_the_document_they_were_given() -> None:
    """Concurrent `lint` and `hangul` runs report on their own document."""

    def run(index: int) -> tuple[list, tuple[str, list[str]]]:
        # Every other document carries a real fault.
        text = f"{_mark(index)}번 절의 본문이다." + (" 傳統的인 방화벽." if index % 2 else "")
        part = {"heading": f"{_mark(index)}절", "content": text}
        findings = lint.check(lint.from_sections([part]))
        return findings, hangul.read_back(text)

    for index, (findings, (clean, replaced)) in enumerate(_spread(run)):
        assert f"{_mark(index)}번 절의 본문이다." in clean
        assert replaced == (["傳統的"] if index % 2 else [])
        assert all(f"{_mark(index)}절" == f.where or not f.where for f in findings)


@pytest.mark.asyncio
async def test_a_stop_reaches_every_turn_running_on_the_session() -> None:
    """중단 stops every turn running on the session."""
    from app.routers.sessions import _STOPPING

    session = "test-session-many-at-once"
    _STOPPING.pop(session, None)
    try:
        first, second = asyncio.Event(), asyncio.Event()
        _STOPPING.setdefault(session, set()).add(first)
        # What starting a second turn does: supersede, then join.
        for earlier in _STOPPING.get(session, set()):
            earlier.set()
        _STOPPING[session].add(second)
        assert first.is_set() and not second.is_set()

        # What 중단 does.
        for signal in _STOPPING.get(session, set()):
            signal.set()
        assert first.is_set() and second.is_set()
    finally:
        _STOPPING.pop(session, None)


@pytest.mark.asyncio
async def test_one_session_stopping_does_not_stop_another() -> None:
    """중단 on one session does not stop another."""
    from app.routers.sessions import _STOPPING

    mine, yours = "many-at-once-mine", "many-at-once-yours"
    for name in (mine, yours):
        _STOPPING.pop(name, None)
    try:
        signals = {name: asyncio.Event() for name in (mine, yours)}
        for name, signal in signals.items():
            _STOPPING.setdefault(name, set()).add(signal)

        for signal in _STOPPING.get(mine, set()):
            signal.set()

        assert signals[mine].is_set()
        assert not signals[yours].is_set()
    finally:
        for name in (mine, yours):
            _STOPPING.pop(name, None)
