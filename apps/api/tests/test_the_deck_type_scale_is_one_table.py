"""One type scale: the panel, the `.pptx` and the `.pdf` read the same table, and the fit
grows a sparse slide as readily as it shrinks a crowded one.
"""

import io
import re
from pathlib import Path

from app.services import deck, deck_export, deck_type

_WEB = Path(__file__).resolve().parents[3] / "apps/web/src/components/slides/typeScale.ts"


def _web_table(name: str) -> dict[str, float]:
    text = _WEB.read_text(encoding="utf-8")
    start = text.index(f"export const {name} = {{")
    block = text[start : text.index("} as const", start)]
    return {
        m.group(1): float(m.group(2)) for m in re.finditer(r"^\s+(\w+): ([\d.]+),", block, re.M)
    }


def test_the_panel_and_the_exporters_read_one_table() -> None:
    """A size changed in one table and not the other would make the file disagree with
    the panel."""
    assert _web_table("TYPE") == {k: float(v) for k, v in deck_type.TYPE.items()}
    assert _web_table("LEADING") == {k: float(v) for k, v in deck_type.LEADING.items()}
    text = _WEB.read_text(encoding="utf-8")
    assert f"export const BULLET_GAP = {deck_type.BULLET_GAP}" in text
    assert f"export const BODY_TOP = {deck_type.BODY_TOP}" in text
    assert f"export const BODY_BOTTOM = {deck_type.BODY_BOTTOM}" in text


def test_a_sparse_list_keeps_its_size_and_a_crowded_one_steps_down() -> None:
    """Nothing grows to fill a slide; a crowded slide takes the next step of the ladder."""
    sparse = {
        "layout": "bullets",
        "title": "왜 문제인가",
        "bullets": ["에이전트는 입력을 명령으로 오인할 수 있음"] * 3,
    }
    assert "textScale" not in deck.auto_fit(dict(sparse))

    crowded = {
        "layout": "bullets",
        "title": "왜 문제인가",
        "bullets": [
            "시스템 권한이 부여된 에이전트는 외부 문서에 숨은 악성 지시를 수행할 위험이 크다"
        ]
        * 9,
    }
    scale = deck.auto_fit(dict(crowded))["textScale"]
    assert scale in deck_type.SCALES and scale < 1
    # Six such lines need one step, nine need more: the ladder is walked, not jumped.
    six = deck.auto_fit({**crowded, "bullets": crowded["bullets"][:6]})["textScale"]
    assert six > scale


def test_the_ladder_is_the_powerpoint_sizes() -> None:
    assert deck_type.STEPS == (22, 18, 16, 14, 12)
    assert deck_type.TYPE["title"] == 32 and deck_type.TYPE["body"] == 22
    assert [round(22 * s) for s in deck_type.SCALES] == [22, 18, 16, 14, 12]


def test_a_title_that_would_wrap_steps_down_to_28_and_no_further() -> None:
    assert deck_type.title_pt("요약") == 32
    assert deck_type.title_pt("문서 검색 시스템 도입 검토 보고서의 주요 결론 정리") == 32
    assert deck_type.title_pt("문서 검색 시스템 도입 검토 보고서의 주요 결론과 권고 사항") == 30
    assert (
        deck_type.title_pt("문서 검색 시스템 도입 검토 보고서의 주요 결론과 권고 사항 정리") == 28
    )
    two_lines = (
        "문서 검색 시스템 도입 검토 보고서의 주요 결론 정리 및 후속 조치 권고 사항과 일정 계획"
    )
    assert deck_type.title_pt(two_lines) == 28
    assert deck_type.lines(two_lines, 28 / deck_type.K, deck_type.TITLE_WIDTH) == 2


def test_a_long_title_takes_its_room_from_the_body() -> None:
    short = {
        "layout": "bullets",
        "title": "요약",
        "bullets": ["짧지 않은 글머리표 한 줄 설명 문장입니다"] * 6,
    }
    long = {
        **short,
        "title": "문서 검색 시스템 도입 검토 보고서의 주요 결론과 후속 조치 권고 사항 정리",
    }
    assert deck.auto_fit(dict(long)).get("textScale", 1.0) <= deck.auto_fit(dict(short)).get(
        "textScale", 1.0
    )


def test_tables_and_cards_size_themselves() -> None:
    table = {"layout": "table", "title": "비교", "rows": [["구분", "값"], ["가", "1"], ["나", "2"]]}
    assert "textScale" not in deck.auto_fit(dict(table))
    cards = {"layout": "cards", "title": "네 가지", "cards": [["이름", "설명"]] * 2}
    assert "textScale" not in deck.auto_fit(dict(cards))
    cover = {"layout": "title", "title": "표지", "bullets": []}
    assert "textScale" not in deck.auto_fit(dict(cover))


def test_a_deck_settles_on_one_body_size() -> None:
    """One crowded slide pulls the deck to 18pt, no lower; it alone goes smaller."""
    line = "시스템 권한이 부여된 에이전트는 외부 문서에 숨은 악성 지시를 수행할 위험이 크다"
    slides = [
        {"layout": "title", "title": "표지"},
        {"layout": "bullets", "title": "하나", "bullets": ["짧은 줄"] * 3},
        {"layout": "bullets", "title": "둘", "bullets": [line] * 9},
        {"layout": "bullets", "title": "셋", "bullets": ["짧은 줄"] * 4},
        {"layout": "closing", "title": "마무리"},
    ]
    for slide in slides:
        deck.auto_fit(slide)
    deck.harmonize(slides)
    eighteen = deck_type.SCALES[1]
    assert slides[1]["textScale"] == eighteen and slides[3]["textScale"] == eighteen
    assert slides[2]["textScale"] < eighteen
    assert "textScale" not in slides[0] and "textScale" not in slides[4]
    assert deck.deck_scale(slides) == eighteen


def test_a_hand_set_scale_is_left_alone() -> None:
    slide = {"layout": "bullets", "title": "요약", "bullets": ["한 줄"], "textScale": 0.9}
    assert deck.auto_fit(dict(slide))["textScale"] == 0.9


def test_columns_follow_their_widest_cell() -> None:
    rows = [
        ["구분", "공격 경로", "대표 예시", "주요 피해"],
        ["직접 주입", "사용자가 직접 입력한 프롬프트", "명령어 위장 입력", "시스템 조작"],
    ]
    shares = deck_type.column_shares(rows)
    assert len(shares) == 4
    assert abs(sum(shares) - 1) < 1e-9
    assert shares[1] > shares[0]
    # One very long cell does not starve the others.
    rows[1][1] = "가" * 200
    assert max(deck_type.column_shares(rows)) < 0.6


def test_a_run_of_bullet_slides_is_broken_up() -> None:
    plan = [{"title": f"{i}", "layout": "bullets"} for i in range(6)]
    plan[0]["layout"] = "title"
    varied = deck.vary_layouts(plan)
    assert [item["layout"] for item in varied] == [
        "title",
        "bullets",
        "bands",
        "bullets",
        "cards",
        "bullets",
    ]
    # Two in a row are fine.
    two = [{"title": "a", "layout": "bullets"}, {"title": "b", "layout": "bullets"}]
    assert [item["layout"] for item in deck.vary_layouts(two)] == ["bullets", "bullets"]


def _every_layout() -> list[dict]:
    pairs = [["30초 QR 갱신", "QR 코드를 30초마다 자동 변경하여 캡처·전송을 무효화합니다."]] * 4
    return [
        {"layout": "title", "title": "스마트 출석관리 시스템", "body": "캡스톤 최종 발표"},
        {"layout": "agenda", "title": "발표 순서", "bullets": [f"{i}장 제목" for i in range(8)]},
        {"layout": "section", "title": "구분", "number": "01"},
        {"layout": "bullets", "title": "왜 문제인가", "bullets": ["한 줄 설명"] * 4},
        {"layout": "two-column", "title": "비교", "bullets": ["항목 설명"] * 6},
        {"layout": "quote", "title": "발표자", "body": "인용문"},
        {"layout": "statement", "title": "핵심 메시지 한 줄", "body": "부연"},
        {
            "layout": "table",
            "title": "표",
            "rows": [["구분", "값", "비고"], ["가", "1", "설명"]] * 2,
        },
        {
            "layout": "metrics",
            "title": "수치",
            "metrics": [["0.3초", "평균 처리"], ["18/20", "만족"]],
        },
        {
            "layout": "big-number",
            "title": "큰 숫자",
            "metrics": [["92%", "차단율"]],
            "body": "설명",
        },
        {
            "layout": "chart",
            "title": "차트",
            "chart": {
                "kind": "bar",
                "categories": ["1월", "2월", "3월"],
                "series": [{"name": "건수", "values": [3, 5, 4]}],
            },
        },
        {"layout": "bands", "title": "항목", "bands": pairs},
        {"layout": "cards", "title": "카드", "cards": pairs},
        {"layout": "steps", "title": "단계", "steps": pairs},
        {"layout": "tiles", "title": "표식", "tiles": [["A", "이름"]] * 3},
        {"layout": "timeline", "title": "연표", "timeline": [["3월", "착수"], ["6월", "완료"]]},
        {"layout": "closing", "title": "마무리", "bullets": ["질문"], "body": "감사합니다"},
    ]


def test_every_layout_exports_at_both_ends_of_the_scale() -> None:
    from pptx import Presentation

    for scale in (0.65, 1.0, 1.25):
        slides = [{**slide, "textScale": scale} for slide in _every_layout()]
        for look in ("editorial", "minimal", "dark"):
            tokens = {"visualStyle": look, "accent": "#1e3a8a"}
            built = Presentation(io.BytesIO(deck_export.to_pptx("확인", slides, tokens=tokens)))
            assert len(built.slides) == len(slides)
            assert deck_export.to_pdf("확인", slides, tokens=tokens).startswith(b"%PDF")


def test_the_title_holds_still_at_every_scale() -> None:
    """The heading is 32pt in the file whatever the body's step, as in the panel."""
    from pptx import Presentation

    def title_size(scale: float) -> float:
        slides = [{"layout": "bullets", "title": "요약", "bullets": ["한 줄"], "textScale": scale}]
        built = Presentation(io.BytesIO(deck_export.to_pptx("확인", slides)))
        return next(
            run.font.size.pt
            for shape in built.slides[0].shapes
            if shape.has_text_frame
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.text == "요약"
        )

    assert title_size(1.25) == title_size(1.0) == title_size(0.8) == deck_type.TYPE["title"]


def test_every_look_the_panel_offers_can_be_exported() -> None:
    """A look listed for the panel has an export twin, a prompt label and a style word."""
    from app.services import design

    assert set(deck_export._LOOKS) == set(design.VISUAL_STYLES)
    assert set(deck._STYLES.values()) == set(design.VISUAL_STYLES)
    assert set(deck._STYLE_LABELS) == set(design.VISUAL_STYLES)
    panel = (
        Path(__file__).resolve().parents[3] / "apps/web/src/components/slides/DeckPanel.tsx"
    ).read_text(encoding="utf-8")
    for look in design.VISUAL_STYLES:
        assert f"\n  {look}: {{ bg:" in panel, f"{look} is missing from the panel's LOOKS"
        assert f"{{ id: '{look}'," in panel, f"{look} is missing from the panel's picker"
    slides = [
        {"layout": "title", "title": "표지"},
        {"layout": "bullets", "title": "본문", "bullets": ["한 줄"]},
        {"layout": "cards", "title": "카드", "cards": [["이름", "설명"]] * 2},
    ]
    for look in ("pastel", "forest", "slate", "paper"):
        tokens = {"visualStyle": look, "accent": "#15803d"}
        assert deck_export.to_pdf("확인", slides, tokens=tokens).startswith(b"%PDF")
        assert deck_export.to_pptx("확인", slides, tokens=tokens)


def test_a_long_deck_opens_with_an_agenda_whatever_the_model_called_it() -> None:
    """The title 「발표 순서」 makes the slide an agenda even when the outline said bullets."""
    plan = [{"title": f"{i}", "layout": "bullets"} for i in range(8)]
    plan[0]["layout"] = "title"
    plan[1]["title"] = "발표 순서"
    assert deck.ensure_agenda(plan)[1]["layout"] == "agenda"
    # No such slide at all: one is put after the cover.
    bare = [{"title": f"{i}", "layout": "bullets"} for i in range(8)]
    bare[0]["layout"] = "title"
    assert deck.ensure_agenda(bare)[1] == {"title": "발표 순서", "layout": "agenda"}
    # A short deck is left alone.
    short = [{"title": "표지", "layout": "title"}, {"title": "발표 순서", "layout": "bullets"}]
    assert deck.ensure_agenda(short)[1]["layout"] == "bullets"
