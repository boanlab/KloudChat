"""Deck table, metrics and chart layouts: cleaning, export, and completion rules."""

from __future__ import annotations

import io
import zipfile

from app.services import deck, deck_export

_ROWS = [
    ["기준", "대안 A", "대안 B"],
    ["초기 비용", "0원", "약 3억"],
    ["도입 기간", "2주", "4개월"],
]


def test_the_model_may_choose_a_table() -> None:
    assert "table" in deck._LAYOUTS
    assert "rows" in deck._TABLE_PROMPT
    assert "머리글" in deck._TABLE_PROMPT


def test_rows_are_cleaned_into_a_rectangle() -> None:
    ragged = deck._clean_rows([["기준", "A", "B"], ["비용", "0원"], ["기간", "2주", "4개월"]])
    # Short rows are padded, not dropped.
    assert ragged == [["기준", "A", "B"], ["비용", "0원", ""], ["기간", "2주", "4개월"]]


def test_a_table_of_one_row_is_not_a_table() -> None:
    assert deck._clean_rows([["기준", "대안 A"]]) == []
    assert deck._clean_rows(None) == []
    assert deck._clean_rows([[], [""]]) == []


def test_a_slide_table_is_cut_to_what_the_back_row_can_read() -> None:
    wide = deck._clean_rows([[f"열{i}" for i in range(9)] for _ in range(9)])
    assert len(wide) == deck._MAX_ROWS
    assert all(len(row) == deck._MAX_COLUMNS for row in wide)
    long = deck._clean_rows([["기준", "A"], ["비용", "가" * 80]])
    assert len(long[1][1]) == deck._MAX_CELL


def test_powerpoint_gets_a_real_table_not_a_list() -> None:
    slides = [
        {"id": "s1", "layout": "title", "title": "제목"},
        {"id": "s2", "layout": "table", "title": "대안 비교", "rows": _ROWS},
    ]
    data = deck_export.to_pptx("제목", slides)
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        slide = archive.read("ppt/slides/slide2.xml").decode()
    # `<a:tbl>` is a native, editable PowerPoint table.
    assert "<a:tbl>" in slide
    for text in ("기준", "대안 A", "약 3억", "4개월"):
        assert f"<a:t>{text}</a:t>" in slide


def test_the_pdf_draws_the_same_table() -> None:
    slides = [{"id": "s1", "layout": "table", "title": "대안 비교", "rows": _ROWS}]
    assert deck_export.to_pdf("제목", slides).startswith(b"%PDF")


def test_the_table_is_drawn_in_the_decks_own_accent() -> None:
    """The .pptx table uses the deck accent, not PowerPoint's banded theme style."""
    slides = [{"id": "s1", "layout": "table", "title": "비교", "accent": "#2b4c7e", "rows": _ROWS}]
    with zipfile.ZipFile(io.BytesIO(deck_export.to_pptx("제목", slides))) as archive:
        slide = archive.read("ppt/slides/slide1.xml").decode()

    # Theme banding must be off or it paints over the explicit fills.
    assert 'firstRow="0"' in slide or "firstRow" not in slide
    assert 'bandRow="0"' in slide or "bandRow" not in slide
    assert "<a:noFill/>" in slide
    assert '<a:srgbClr val="2B4C7E"/>' in slide
    assert "<a:lnB " in slide


def test_the_head_rule_sits_where_the_schema_wants_it() -> None:
    """Inside `a:tcPr` the border comes before the fill, as the schema orders them."""
    import re

    slides = [{"id": "s1", "layout": "table", "title": "비교", "rows": _ROWS}]
    with zipfile.ZipFile(io.BytesIO(deck_export.to_pptx("제목", slides))) as archive:
        slide = archive.read("ppt/slides/slide1.xml").decode()
    for properties in re.findall(r"<a:tcPr\b[^>]*>(.*?)</a:tcPr>", slide, re.S):
        if "<a:lnB" in properties and "<a:noFill/>" in properties:
            assert properties.index("<a:lnB") < properties.index("<a:noFill/>")


# ── 강조 수치 ──────────────────────────────────────────────────────────

_METRICS = [["32%", "오탐 감소"], ["1.4초", "평균 응답"], ["99.2%", "가용성"]]


def test_the_model_may_choose_metrics() -> None:
    assert "metrics" in deck._LAYOUTS
    assert "지어낸 수치" in deck._METRICS_PROMPT
    assert "bullets 로 답하라" in deck._METRICS_PROMPT


def test_a_figure_with_no_label_is_dropped() -> None:
    assert deck._clean_metrics([["32%"], ["1.4초", "평균 응답"], ["", "가용성"]]) == [
        ["1.4초", "평균 응답"]
    ]
    assert deck._clean_metrics(None) == []


def test_metrics_are_cut_to_what_fits_across_a_slide() -> None:
    many = deck._clean_metrics([[f"{n}%", f"이름{n}"] for n in range(9)])
    assert len(many) == deck._MAX_METRICS
    long = deck._clean_metrics([["가" * 40, "나" * 40]])
    assert len(long[0][0]) == deck._MAX_VALUE
    assert len(long[0][1]) == deck._MAX_LABEL


def test_powerpoint_sets_the_figures_large() -> None:
    """Metric figures are set at 44pt with labels well below 20pt."""
    import re

    slides = [{"id": "s1", "layout": "metrics", "title": "성과", "metrics": _METRICS}]
    with zipfile.ZipFile(io.BytesIO(deck_export.to_pptx("제목", slides))) as archive:
        slide = archive.read("ppt/slides/slide1.xml").decode()
    for text in ("32%", "오탐 감소", "99.2%"):
        assert f"<a:t>{text}</a:t>" in slide
    sizes = {int(size) for size in re.findall(r'sz="(\d+)"', slide)}
    assert 4400 in sizes
    assert min(sizes) < 2000


def test_the_pdf_draws_the_same_figures() -> None:
    slides = [{"id": "s1", "layout": "metrics", "title": "성과", "metrics": _METRICS}]
    assert deck_export.to_pdf("제목", slides).startswith(b"%PDF")


# ── 차트 ───────────────────────────────────────────────────────────────

_CHART = {
    "kind": "bar",
    "unit": "건",
    "categories": ["1분기", "2분기", "3분기", "4분기"],
    "series": [{"name": "처리 건수", "values": [120, 210, 380, 460]}],
}


def test_the_model_may_choose_a_chart() -> None:
    assert "chart" in deck._LAYOUTS
    assert "지어낸 수치" in deck._CHART_PROMPT
    assert "bullets 로 답하라" in deck._CHART_PROMPT


def test_a_series_shorter_than_its_categories_is_trimmed_not_padded() -> None:
    """Categories and values are both cut to the length they share."""
    trimmed = deck._clean_chart(
        {
            "kind": "bar",
            "categories": ["1분기", "2분기", "3분기", "4분기"],
            "series": [{"name": "건수", "values": [120, 210, 380]}],
        }
    )
    assert trimmed is not None
    assert trimmed["categories"] == ["1분기", "2분기", "3분기"]
    assert trimmed["series"][0]["values"] == [120, 210, 380]


def test_a_chart_that_cannot_be_drawn_is_refused() -> None:
    assert deck._clean_chart(None) is None
    assert deck._clean_chart({"categories": ["가", "나"], "series": []}) is None
    assert deck._clean_chart({"categories": ["가"], "series": [{"values": [1]}]}) is None
    # A non-numeric value truncates the series there; skipping it would misalign labels.
    partial = deck._clean_chart(
        {"categories": ["가", "나", "다", "라"], "series": [{"values": [1, 2, "미정", 4]}]}
    )
    assert partial["categories"] == ["가", "나"]
    assert partial["series"][0]["values"] == [1.0, 2.0]
    assert (
        deck._clean_chart(
            {"categories": ["가", "나", "다"], "series": [{"values": [1, "미정", 3]}]}
        )
        is None
    )


def test_powerpoint_gets_a_chart_it_can_edit() -> None:
    """The .pptx carries a native chart part with an embedded workbook and a zero floor."""
    slides = [{"id": "s1", "layout": "chart", "title": "처리 추이", "chart": _CHART}]
    with zipfile.ZipFile(io.BytesIO(deck_export.to_pptx("제목", slides))) as archive:
        parts = archive.namelist()
        assert any(n.startswith("ppt/charts/chart") for n in parts), parts
        assert any("embeddings" in n for n in parts), parts
        chart = next(archive.read(n).decode() for n in parts if n.startswith("ppt/charts/chart"))

    for label in ("1분기", "4분기", "처리 건수"):
        assert label in chart
    assert "<c:min val=\"0\"/>" in chart or "<c:min val=\"0.0\"/>" in chart


def test_the_pdf_draws_the_same_chart() -> None:
    slides = [{"id": "s1", "layout": "chart", "title": "처리 추이", "chart": _CHART}]
    assert deck_export.to_pdf("제목", slides).startswith(b"%PDF")
    line = [
        {"id": "s1", "layout": "chart", "title": "추이", "chart": dict(_CHART, kind="line")}
    ]
    assert deck_export.to_pdf("제목", line).startswith(b"%PDF")


def test_an_unusable_chart_does_not_break_the_export() -> None:
    broken = [
        {"id": "s1", "layout": "chart", "title": "x", "chart": {"categories": [], "series": []}}
    ]
    assert deck_export.to_pptx("제목", broken)
    assert deck_export.to_pdf("제목", broken).startswith(b"%PDF")


def test_the_printed_scale_uses_numbers_a_reader_recognises() -> None:
    """The PDF axis ceiling is a round number whose quarters read as a scale."""
    from app.services.deck_export import _nice_ceiling

    assert _nice_ceiling(460) == 500
    assert _nice_ceiling(8.1) == 10
    assert _nice_ceiling(0.9) == 1
    assert _nice_ceiling(1200) == 2000

    from app.services.deck_export import _tick_label

    for highest in (3, 47, 460, 8.1, 12345):
        top = _nice_ceiling(highest)
        assert highest <= top <= highest * 2.5
        assert all(len(_tick_label(top * n / 4)) <= 7 for n in range(5))


def test_the_printed_chart_names_its_lines() -> None:
    """A multi-series PDF line chart draws a legend naming each series."""
    from app.services import deck_export as export

    drawn: list[str] = []

    class _Spy:
        def __getattr__(self, name):
            def record(*args, **kwargs):
                if name == "drawString":
                    drawn.append(str(args[-1]))
                return _Spy()

            return record

    export._pdf_chart(
        _Spy(),
        {
            "kind": "line",
            "unit": "초",
            "categories": ["1월", "2월", "3월"],
            "series": [("평균", [3.0, 2.0, 1.0]), ("95분위", [8.0, 6.0, 4.0])],
        },
        accent=(0.1, 0.3, 0.5),
        muted=(0.4, 0.4, 0.4),
        top=400,
        width=600,
        font="Helvetica",
    )
    assert "평균" in drawn and "95분위" in drawn


# ── 완성 판정 ──────────────────────────────────────────────────────────

def test_a_slide_of_the_new_kinds_counts_as_written() -> None:
    """A slide whose content is rows, metrics or a chart counts as written."""
    for slide in (
        {"layout": "table", "title": "비교", "rows": _ROWS},
        {"layout": "metrics", "title": "성과", "metrics": _METRICS},
        {"layout": "chart", "title": "추이", "chart": _CHART},
    ):
        assert deck.has_content(slide), slide["layout"]
        assert deck.filled([slide]) == [slide], slide["layout"]


def test_an_empty_slide_still_counts_as_unwritten() -> None:
    assert not deck.has_content({"layout": "bullets", "title": "빈 장"})
    assert not deck.has_content({"layout": "bullets", "title": "빈 장", "body": "   "})
    assert deck.filled([{"layout": "bullets", "title": "빈 장"}]) == []
    assert deck.has_content({"layout": "title", "title": "표지"})


def test_a_slide_that_came_back_wrong_still_gets_something() -> None:
    """Bullets are salvaged from a misshapen model answer, but not from notes or structure."""
    assert deck._salvaged_bullets({"points": ["가상환경 격리", "의존성 고정"]}) == [
        "가상환경 격리",
        "의존성 고정",
    ]
    # A paragraph is split at sentence ends.
    assert len(deck._salvaged_bullets({"content": "첫 문장이다. 둘째 문장이다."})) == 2
    assert deck._salvaged_bullets({"notes": "말로 할 이야기", "layout": "bullets"}) == []
    assert deck._salvaged_bullets({"data": {"a": 1}, "count": 3}) == []


def test_the_salvage_only_runs_when_there_is_nothing_else() -> None:
    """Salvage runs only for a slide with no content of its own."""
    import inspect

    source = inspect.getsource(deck._write_slides)
    assert "if not has_content(slide):" in source
    assert "_salvaged_bullets(data)" in source


def test_a_slide_never_comes_out_blank() -> None:
    """An unwritable slide gets the UNWRITTEN marker as body, after salvage has failed."""
    import inspect

    source = inspect.getsource(deck._write_slides)
    marker = 'slide["body"] = UNWRITTEN'
    # Set both when the call throws and when nothing usable came back.
    assert source.count(marker) == 2
    assert source.index("_salvaged_bullets(data)") < source.rindex(marker)
    # The marker counts as content.
    assert deck.has_content({"layout": "bullets", "body": deck.UNWRITTEN})


def test_the_sentence_says_so_on_screen_and_not_in_the_file() -> None:
    """An UNWRITTEN-only slide is dropped from the exported file, not left blank."""
    import io

    from pptx import Presentation

    from app.services import deck_export

    slides = [
        {"layout": "title", "title": "덱", "body": "부제"},
        {"layout": "bullets", "title": "쓴 장", "bullets": ["내용"]},
        {"layout": "bullets", "title": "못 쓴 장", "body": deck.UNWRITTEN},
    ]
    written = Presentation(io.BytesIO(deck_export.to_pptx("덱", slides)))
    assert len(written.slides._sldIdLst) == 2
    words = " ".join(
        shape.text_frame.text
        for slide in written.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert deck.UNWRITTEN not in words
    assert "쓴 장" in words


def test_a_slide_that_only_lost_its_prose_is_kept() -> None:
    """A slide marked UNWRITTEN but carrying a table is still exported."""
    import io

    from pptx import Presentation

    from app.services import deck_export

    slides = [
        {"layout": "title", "title": "덱"},
        {
            "layout": "table",
            "title": "표는 나왔다",
            "body": deck.UNWRITTEN,
            "rows": [["가", "나"], ["1", "2"]],
        },
    ]
    written = Presentation(io.BytesIO(deck_export.to_pptx("덱", slides)))
    assert len(written.slides._sldIdLst) == 2
