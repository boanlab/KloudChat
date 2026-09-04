"""Exported .pptx slides carry `p:ph` placeholders so PowerPoint's outline pane reads them."""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from app.services import deck_export

_SLIDES = [
    {"layout": "title", "title": "전교생 AI기초 교육 의무화", "body": "2026 개편안"},
    {
        "layout": "bullets",
        "title": "현행의 문제",
        "bullets": ["이수체계가 학과마다 다르다", "실습 비중이 낮다"],
    },
    {"layout": "quote", "title": "교무처장", "body": "전교생이 같은 출발선에 서야 한다"},
    {"layout": "table", "title": "비교", "rows": [["구분", "기존", "개편"], ["학점", "6", "9"]]},
    {"layout": "metrics", "title": "규모", "metrics": [["9학점", "필수"], ["17명", "AI-PD"]]},
]


def _placeholders(slide) -> list[tuple[str, str]]:
    """`(type, text)` for every shape on the slide that is outline text."""
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        found = shape.element.nvSpPr.nvPr.find(qn("p:ph"))
        if found is not None:
            out.append((found.get("type"), shape.text_frame.text))
    return out


def _opened(slides: list[dict]) -> Presentation:
    return Presentation(io.BytesIO(deck_export.to_pptx("AI 교육", slides)))


def test_every_slide_says_what_it_is() -> None:
    """Every slide has at least one placeholder shape."""
    deck = _opened(_SLIDES)
    titled = [s for s in deck.slides if _placeholders(s)]
    assert len(titled) == len(_SLIDES)


def test_the_cover_is_a_centred_title_and_the_rest_are_titles() -> None:
    """The cover is `ctrTitle`; every other slide is `title`."""
    kinds = [_placeholders(s)[0][0] for s in _opened(_SLIDES).slides]
    assert kinds[0] == "ctrTitle"
    assert set(kinds[1:]) == {"title"}


def test_the_cover_placeholder_keeps_the_drawn_left_alignment() -> None:
    """`ctrTitle` identifies the cover without restyling its visible textbox."""
    cover = _opened(_SLIDES).slides[0]
    title = next(
        shape for shape in cover.shapes
        if shape.has_text_frame
        and "전교생 AI기초 교육 의무화" in shape.text_frame.text
        and shape.element.nvSpPr.nvPr.find(qn("p:ph")) is None
    )
    assert title.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT
    assert title.text_frame.paragraphs[1].alignment == PP_ALIGN.LEFT
    semantic = next(
        shape for shape in cover.shapes
        if shape.has_text_frame and shape.element.nvSpPr.nvPr.find(qn("p:ph")) is not None
    )
    assert semantic.element.nvSpPr.cNvPr.get("hidden") == "1"


def test_the_bullets_are_the_body_so_the_outline_holds_the_content() -> None:
    """Bullets are the `body` placeholder."""
    bullets = _opened(_SLIDES).slides[1]
    kinds = dict(_placeholders(bullets))
    assert "이수체계가 학과마다 다르다" in kinds["body"]
    assert "실습 비중이 낮다" in kinds["body"]


def test_the_words_still_read_the_way_they_were_drawn() -> None:
    """The title placeholder text is exactly the slide's title."""
    slide = _opened(_SLIDES).slides[1]
    title = next(text for kind, text in _placeholders(slide) if kind == "title")
    assert title == "현행의 문제"
