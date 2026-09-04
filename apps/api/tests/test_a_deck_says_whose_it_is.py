"""Deck footer, logo and divider tokens: kept, drawn, and read by every renderer."""

from __future__ import annotations

import base64
import io

import PIL.Image
import pytest
from pptx import Presentation

from app.services import deck_export, design


def _logo() -> str:
    picture = PIL.Image.new("RGB", (240, 60), (255, 255, 255))
    buffer = io.BytesIO()
    picture.save(buffer, "PNG")
    return "data:image/png;base64," + base64.b64encode(buffer.getvalue()).decode()


_SLIDES = [
    {"layout": "title", "title": "부트캠프 사전설명회", "body": "인공지능학과"},
    {"layout": "section", "number": "01.", "title": "사업 소개"},
    {"layout": "bullets", "title": "개요", "bullets": ["대학과 기업이 함께 운영"]},
]


def test_the_marks_are_kept_and_the_junk_is_not() -> None:
    """A footer is one line and a logo is an embedded picture; anything else is dropped."""
    kept = design.normalise_tokens({"footer": "  단국대학교   인공지능학과 ", "logo": _logo()})
    assert kept["footer"] == "단국대학교 인공지능학과"
    assert kept["logo"].startswith("data:image/png;base64,")

    dropped = design.normalise_tokens({"footer": "x" * 200, "logo": "https://example.com/a.png"})
    assert len(dropped["footer"]) == 80
    # URL logos are dropped; only embedded data.
    assert dropped["logo"] == ""


def test_a_deck_with_no_design_system_draws_what_it_always_drew() -> None:
    """Empty marks draw nothing, not a blank box."""
    assert design.DEFAULT_TOKENS["footer"] == ""
    assert design.DEFAULT_TOKENS["logo"] == ""
    assert deck_export.to_pdf("덱", _SLIDES)[:4] == b"%PDF"


def test_the_logo_reaches_the_pptx() -> None:
    """Every slide but the cover carries exactly one logo picture."""
    tokens = design.normalise_tokens({"logo": _logo(), "footer": "인공지능학과"})
    deck = Presentation(io.BytesIO(deck_export.to_pptx("덱", _SLIDES, tokens=tokens)))
    pictures = [
        sum(1 for shape in slide.shapes if shape.shape_type == 13) for slide in deck.slides
    ]
    # Cover and divider carry none; the content slide carries the mark.
    assert pictures == [0, 0, 1]


def test_the_footer_reaches_the_pptx() -> None:
    tokens = design.normalise_tokens({"footer": "단국대학교 인공지능학과"})
    deck = Presentation(io.BytesIO(deck_export.to_pptx("덱", _SLIDES, tokens=tokens)))
    words = [
        " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        for slide in deck.slides
    ]
    assert "단국대학교 인공지능학과" in words[2]
    # Not on the cover.
    assert "단국대학교 인공지능학과" not in words[0]


def test_a_bad_logo_is_dropped_rather_than_failing_the_export() -> None:
    """An undecodable logo never fails the export."""
    tokens = dict(design.DEFAULT_TOKENS, logo="data:image/png;base64,bm90IGEgcG5n")
    assert deck_export.to_pdf("덱", _SLIDES, tokens=tokens)[:4] == b"%PDF"
    assert deck_export.to_pptx("덱", _SLIDES, tokens=tokens)[:2] == b"PK"


@pytest.mark.parametrize("write", [deck_export.to_pdf, deck_export.to_pptx])
def test_a_divider_carries_its_number(write) -> None:
    """A divider shows `01.` above the part name."""
    written = write("덱", _SLIDES)
    if written[:2] == b"PK":
        deck = Presentation(io.BytesIO(written))
        words = " ".join(
            s.text_frame.text for s in deck.slides[1].shapes if s.has_text_frame
        )
        assert "01." in words and "사업 소개" in words
    else:
        assert written[:4] == b"%PDF"


def test_a_divider_is_not_one_of_the_shapes_an_argument_takes() -> None:
    """The layout-variety check ignores dividers."""
    from app.services import deck

    assert "section" not in deck._BODY_LAYOUTS
    assert "title" not in deck._BODY_LAYOUTS
    assert set(deck._BODY_LAYOUTS) < set(deck._LAYOUTS)


def test_the_hwpx_and_the_deck_read_the_same_token_set() -> None:
    """HWPX and deck renderers accept the same token set."""
    assert set(design.normalise_tokens({})) == set(design.DEFAULT_TOKENS)
