"""Paired deck layouts (bands, tiles, timeline, steps, cards) and the structural slides."""

from __future__ import annotations

import io

import pytest
from pptx import Presentation

from app.services import deck, deck_export

_SLIDES = {
    "bands": [["미션", "차세대 AX 인재를 양성한다"], ["배경", "인재 격차가 벌어지고 있다"]],
    "tiles": [["P", "Physical AI"], ["H", "Human-centered AI"]],
    "timeline": [["2024.05", "1기 개설"], ["2025.03", "2기 개설"]],
    "steps": [["접수", "신청서를 낸다"], ["심사", "서류와 발표로 평가한다"]],
    "cards": [["교육", "전교생 필수 과목을 연다"], ["연구", "융합 과제를 지원한다"]],
}


@pytest.mark.parametrize("layout", sorted(_SLIDES))
def test_each_shape_reaches_both_files(layout: str) -> None:
    slide = {"layout": layout, "title": "장", layout: _SLIDES[layout]}
    assert deck_export.to_pdf("덱", [slide])[:4] == b"%PDF"

    written = Presentation(io.BytesIO(deck_export.to_pptx("덱", [slide])))
    words = " ".join(
        shape.text_frame.text for shape in written.slides[0].shapes if shape.has_text_frame
    )
    for left, right in _SLIDES[layout]:
        assert left in words
        assert right in words


@pytest.mark.parametrize("layout", sorted(_SLIDES))
def test_a_half_empty_pair_is_dropped(layout: str) -> None:
    """A pair missing either half is dropped."""
    pairs = deck._clean_pairs([["이름", "내용"], ["이름만"], ["", "내용만"], ["둘째", ""]], layout)
    assert pairs == [["이름", "내용"]]


@pytest.mark.parametrize("layout", sorted(_SLIDES))
def test_an_empty_shape_becomes_bullets_rather_than_a_blank_slide(layout: str) -> None:
    """An empty or non-list value cleans to an empty list."""
    assert deck._clean_pairs([], layout) == []
    assert deck._clean_pairs("not a list", layout) == []


def test_the_halves_are_bounded_differently_because_they_hold_different_things() -> None:
    """Left halves are bounded as names, right halves as sentences."""
    long_pair = [["가" * 200, "나" * 200]]
    left, right = deck._clean_pairs(long_pair, "bands")[0]
    assert len(left) == 10
    assert len(right) == 90
    # A tile's mark is at most 4 characters.
    assert len(deck._clean_pairs(long_pair, "tiles")[0][0]) == 4


@pytest.mark.parametrize("layout", sorted(_SLIDES))
def test_the_writer_is_told_what_each_half_holds(layout: str) -> None:
    """Each paired layout's prompt names the layout and takes heading and outline."""
    prompt = deck._PROMPTS[layout]
    assert layout in prompt
    assert "{heading}" in prompt and "{outline}" in prompt


@pytest.mark.parametrize("layout", sorted(_SLIDES))
def test_they_are_shapes_an_argument_takes(layout: str) -> None:
    """Paired layouts are body layouts and none is numeric."""
    assert layout in deck._BODY_LAYOUTS
    assert layout not in deck._NUMERIC_LAYOUTS


def test_an_artifact_written_before_these_existed_still_draws() -> None:
    """A paired layout with only `bullets` still renders in both exporters."""
    stale = {"layout": "bands", "title": "장", "bullets": ["옛 문서"]}
    assert deck_export.to_pdf("덱", [stale])[:4] == b"%PDF"
    assert deck_export.to_pptx("덱", [stale])[:2] == b"PK"


def test_a_filled_shape_is_content() -> None:
    """Every paired layout is in `_CONTENT_FIELDS` and a filled one counts as content."""
    for layout in sorted(_SLIDES):
        assert layout in deck._CONTENT_FIELDS
        assert deck.has_content({"layout": layout, layout: _SLIDES[layout]})


def test_the_marker_is_not_written_over_a_slide_that_has_something_on_it() -> None:
    """Filled bands count as content even without body prose."""
    slide = {"layout": "bands", "title": "혜택", "bands": _SLIDES["bands"]}
    assert deck.has_content(slide)


def test_the_structural_slides_reach_both_files() -> None:
    """Agenda, statement, big-number and closing slides render in both exporters."""
    slides = [
        {"layout": "title", "title": "덱"},
        {"layout": "agenda", "title": "목차", "bullets": ["배경", "방법", "일정"]},
        {"layout": "statement", "title": "전교생, AI 기초부터", "body": "같은 출발선에서"},
        {
            "layout": "big-number",
            "title": "성과",
            "metrics": [["32%", "오탐 감소"]],
            "body": "첫 달",
        },
        {"layout": "closing", "title": "마무리", "bullets": ["기억할 것"], "body": "감사합니다"},
    ]
    assert deck_export.to_pdf("덱", slides)[:4] == b"%PDF"
    written = Presentation(io.BytesIO(deck_export.to_pptx("덱", slides)))
    texts = [
        " ".join(s.text_frame.text for s in page.shapes if s.has_text_frame)
        for page in written.slides
    ]
    assert "01" in texts[1] and "배경" in texts[1]
    assert "전교생, AI 기초부터" in texts[2]
    assert "32%" in texts[3] and "오탐 감소" in texts[3]
    assert "기억할 것" in texts[4] and "감사합니다" in texts[4]


def test_the_agenda_is_the_outline_read_back() -> None:
    from app.services import deck

    plan = [
        {"layout": "title", "title": "덱"},
        {"layout": "agenda", "title": "목차"},
        {"layout": "bullets", "title": "배경"},
        {"layout": "steps", "title": "절차"},
        {"layout": "closing", "title": "마무리"},
    ]
    assert deck._agenda_lines(plan) == ["배경", "절차"]
    with_parts = [
        *plan[:2],
        {"layout": "section", "title": "1부"},
        plan[2],
        {"layout": "section", "title": "2부"},
        plan[3],
        plan[4],
    ]
    assert deck._agenda_lines(with_parts) == ["1부", "2부"]
