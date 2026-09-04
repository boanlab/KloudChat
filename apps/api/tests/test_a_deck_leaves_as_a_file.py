"""Deck 서식 reach the `.pptx`: master and layouts carry the design, read from the built file."""

from __future__ import annotations

import io
import pathlib

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from app.services import design_templates

_THEME = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


def _decks() -> list:
    return [row for row in design_templates.all_templates() if row.kind == "deck"]


def test_every_deck_template_ships_a_form() -> None:
    for row in _decks():
        assert row.form_file, f"{row.id} 에 form.pptx 가 없습니다"
        assert pathlib.Path(row.form_file).is_file()
        assert row.form_file.endswith(".pptx"), f"{row.id}: 덱의 양식은 .pptx 여야 합니다"


def test_the_form_is_made_of_layouts_and_not_of_drawings() -> None:
    """모든 글자가 자리표시자 안에 있다."""
    for row in _decks():
        deck = Presentation(row.form_file)
        assert len(deck.slides) > 0, f"{row.id}: 양식에 슬라이드가 없습니다"
        for index, slide in enumerate(deck.slides, start=1):
            drawn = [shape.name for shape in slide.shapes if not shape.is_placeholder]
            assert not drawn, f"{row.id} {index}장: 직접 그린 도형 {drawn}"
            assert slide.slide_layout.name, f"{row.id} {index}장: 레이아웃이 없습니다"


def test_the_design_is_on_the_master_rather_than_on_the_runs() -> None:
    """글꼴도 색도 크기도 마스터가 들고 있다."""
    #: First-level sizes seen across the catalogue, per style list.
    scales: dict[str, set[int]] = {}

    for row in _decks():
        deck = Presentation(row.form_file)

        for slide in deck.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        assert run.font.size is None, f"{row.id}: 글자 크기를 직접 지정했습니다"
                        assert run.font.name is None, f"{row.id}: 글꼴을 직접 지정했습니다"
                        assert run.font.color.type is None, f"{row.id}: 색을 직접 지정했습니다"

        master = deck.slide_masters[0]
        theme = etree.fromstring(master.part.part_related_by(_THEME).blob)
        elements = theme.find(qn("a:themeElements"))

        fonts = elements.find(qn("a:fontScheme"))
        for tag in ("a:majorFont", "a:minorFont"):
            block = fonts.find(qn(tag))
            latin = block.find(qn("a:latin")).get("typeface")
            east = block.find(qn("a:ea")).get("typeface")
            assert latin and latin != "Calibri Light", f"{row.id}: {tag} 가 기본값입니다"
            # A Hangul run reads `ea`, not `latin`.
            assert east == latin, f"{row.id}: {tag} 의 한글 글꼴이 라틴과 다릅니다"

        styles = master.element.find(qn("p:txStyles"))
        assert styles is not None, f"{row.id}: 마스터에 텍스트 스타일이 없습니다"
        for tag in ("p:titleStyle", "p:bodyStyle"):
            block = styles.find(qn(tag))
            levels = []
            for level in list(block):
                run = level.find(qn("a:defRPr"))
                assert run is not None, f"{row.id}: {tag} 에 defRPr 없는 수준이 있습니다"
                size = run.get("sz")
                assert size, f"{row.id}: {tag} 에 크기 없는 수준이 있습니다"
                levels.append(int(size))
            # Outline levels step down.
            assert levels == sorted(levels, reverse=True), (
                f"{row.id}: {tag} 크기가 내려가지 않습니다"
            )
            scales.setdefault(tag, set()).add(levels[0])


    # The 서식 declare different scales; identical numbers mean none took.
    for tag, seen in scales.items():
        assert len(seen) > 1, (
            f"{tag}: 모든 덱이 같은 크기 {seen} 입니다 — 서식이 반영되지 않았습니다"
        )


def test_the_template_ships_without_slides() -> None:
    """`template.pptx` 는 마스터와 레이아웃뿐이다."""
    for row in _decks():
        base = pathlib.Path(row.form_file).with_name("template.pptx")
        assert base.is_file(), f"{row.id} 에 template.pptx 가 없습니다"
        assert len(Presentation(str(base)).slides) == 0, f"{row.id}: 템플릿에 슬라이드가 있습니다"


def test_the_contents_use_the_wide_canvas() -> None:
    """16:9 캔버스의 자리표시자가 4:3 좌표에 남아 있지 않다."""
    for row in _decks():
        deck = Presentation(row.form_file)
        assert deck.slide_width / deck.slide_height > 1.7, f"{row.id}: 16:9 가 아닙니다"

        edges = [
            (shape.left + shape.width) / deck.slide_width
            for slide in deck.slides
            for shape in slide.placeholders
            if shape.left is not None and shape.width is not None
        ]
        assert edges, f"{row.id}: 잴 자리표시자가 없습니다"
        # 4:3 geometry on a 16:9 canvas stops at 8,686,800 of 12,192,000 EMU (71%).
        assert max(edges) > 0.8, (
            f"{row.id}: 가장 오른쪽 자리표시자가 {max(edges):.0%} 에서 멈춥니다 — "
            "마스터와 레이아웃이 4:3 좌표에 남아 있습니다"
        )
        # Scaled once: an inheriting placeholder already reports the widened value.
        assert max(edges) <= 1.0, (
            f"{row.id}: 자리표시자가 슬라이드 밖 {max(edges):.0%} 까지 나갑니다 — "
            "물려받은 값을 두 번 키웠을 때 이렇게 됩니다"
        )

        # An inheriting placeholder has no `a:xfrm`; setting one value zeroes
        # the other half of the pair (`left` alone gives `y="0"`).
        for slide in deck.slides:
            for shape in slide.placeholders:
                assert shape.height, f"{row.id}: '{shape.name}' 높이가 0 입니다"
                assert shape.top + shape.height <= deck.slide_height, (
                    f"{row.id}: '{shape.name}' 이 슬라이드 아래로 넘칩니다"
                )


def test_the_export_is_built_on_the_template_it_was_written_in() -> None:
    """덱도 제 서식의 파일 위에 지어진다."""
    import io

    from app.services import deck_export

    slides = [
        {"layout": "title", "title": "제목", "body": "부제"},
        {"layout": "bullets", "title": "본문", "bullets": ["가", "나"]},
    ]

    def major(raw: bytes) -> str:
        deck = Presentation(io.BytesIO(raw))
        theme = etree.fromstring(deck.slide_masters[0].part.part_related_by(_THEME).blob)
        fonts = theme.find(qn("a:themeElements")).find(qn("a:fontScheme"))
        return fonts.find(qn("a:majorFont")).find(qn("a:latin")).get("typeface")

    # No 서식: PowerPoint's own face.
    assert major(deck_export.to_pptx("제목", slides)) == "Calibri"

    for row in _decks():
        assert row.pptx_template, f"{row.id} 에 template.pptx 가 없습니다"
        built = deck_export.to_pptx("제목", slides, template=row.pptx_template)
        assert major(built) != "Calibri", f"{row.id}: 서식이 파일에 실리지 않았습니다"


def test_a_slide_can_be_set_larger_or_smaller_and_the_file_follows() -> None:
    """글자 크기 on one slide is honoured by the `.pptx`, scaled once in `paint`."""
    from pptx import Presentation

    from app.services import deck_export

    body = [{"id": "s1", "layout": "bullets", "title": "비밀번호", "bullets": ["12자 이상"]}]

    def sizes(scale: float | None) -> list[float]:
        slides = [{**body[0], **({"textScale": scale} if scale else {})}]
        deck = Presentation(io.BytesIO(deck_export.to_pptx("확인", slides)))
        return sorted(
            run.font.size.pt
            for slide in deck.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
            if run.font.size
        )

    assert sizes(1.25)[-1] > sizes(None)[-1] > sizes(0.8)[-1]


def test_selected_text_formatting_survives_into_powerpoint() -> None:
    """Inline emphasis reaches the downloaded file."""
    from app.services import deck_export

    slides = [{
        "id": "s1",
        "layout": "bullets",
        "title": "선택 서식",
        "bullets": ["중요 내용"],
        "richText": {
            "bullets.0": '<b><span style="font-size:1.35em;color:#c02020">중요</span></b> 내용'
        },
    }]
    built = Presentation(io.BytesIO(deck_export.to_pptx("확인", slides)))
    runs = [
        run
        for shape in built.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if "중요" in run.text
    ]
    assert len(runs) == 1
    assert runs[0].font.bold is True
    assert runs[0].font.size.pt > 18
    assert str(runs[0].font.color.rgb) == "C02020"


def test_an_absurd_type_size_cannot_be_stored_into_the_file() -> None:
    """`textScale` on a PATCHed artifact is clamped by the exporter."""
    from app.services import deck_export

    assert deck_export._typescale({"textScale": 40}) == 2.0
    assert deck_export._typescale({"textScale": 0.1}) == 0.5
    assert deck_export._typescale({"textScale": "크게"}) == 1.0
    assert deck_export._typescale({}) == 1.0
    # `0` is how clients serialise an absent value; read as unset, not clamped.
    assert deck_export._typescale({"textScale": 0}) == 1.0
